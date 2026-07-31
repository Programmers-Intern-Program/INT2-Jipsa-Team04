"""Issue #159 자원·한계 측정을 위한 결정적 다중 형식 문서를 생성한다.

생성 문서는 운영 데이터나 사용자 파일을 사용하지 않는다. 같은 계획과 File_IDX를
사용하면 동일한 텍스트와 이미지 패턴을 생성하여 실행 간 비교 가능성을 높인다.
문서 크기, 텍스트 단위 수와 이미지 수만 단계적으로 늘리며 서비스의 청킹·OCR·임베딩
정책 자체는 변경하지 않는다.
"""

from __future__ import annotations

import hashlib
import json
import textwrap
from collections.abc import Iterable
from pathlib import Path
from typing import Final

import fitz  # type: ignore[import-untyped]
from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import Workbook  # type: ignore[import-untyped]
from openpyxl.drawing.image import Image as OpenpyxlImage  # type: ignore[import-untyped]
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches as PptxInches

from jipsa_rag_benchmark.models import (
    BenchmarkPlan,
    FixtureProfile,
    GeneratedFixture,
    SupportedFormat,
)

_CONTENT_TYPE_BY_FORMAT: Final[dict[SupportedFormat, str]] = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "txt": "text/plain",
}

_EXTENSION_BY_FORMAT: Final[dict[SupportedFormat, str]] = {
    file_type: f".{file_type}" for file_type in _CONTENT_TYPE_BY_FORMAT
}


class FixtureFactory:
    """계획에 선언된 형식·프로필 조합을 실제 문서로 생성한다."""

    def __init__(self, *, output_directory: Path, file_idx_start: int) -> None:
        self._output_directory = output_directory
        self._next_file_idx = file_idx_start
        self._image_directory = output_directory / "_generated_images"

    def generate(self, plan: BenchmarkPlan) -> tuple[GeneratedFixture, ...]:
        """중복 조합을 제거하고 계획 순서대로 Fixture를 생성한다."""

        self._output_directory.mkdir(parents=True, exist_ok=True)
        self._image_directory.mkdir(parents=True, exist_ok=True)

        profiles = plan.profiles_by_name
        generated: list[GeneratedFixture] = []
        seen_combinations: set[tuple[SupportedFormat, str]] = set()

        for matrix_entry in plan.fixture_matrix:
            for file_type in matrix_entry.formats:
                for profile_name in matrix_entry.profiles:
                    combination = (file_type, profile_name)
                    if combination in seen_combinations:
                        continue
                    seen_combinations.add(combination)

                    profile = profiles[profile_name]
                    fixture = self._generate_one(
                        group=matrix_entry.group,
                        file_type=file_type,
                        profile=profile,
                    )
                    generated.append(fixture)

        manifest_path = self._output_directory / "fixtures_manifest.json"
        manifest_path.write_text(
            json.dumps(
                {
                    "schema_version": 1,
                    "fixtures": [fixture.to_public_dict() for fixture in generated],
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        return tuple(generated)

    def allocate_clone(
        self,
        source: GeneratedFixture,
        *,
        purpose: str,
        ordinal: int,
    ) -> GeneratedFixture:
        """동시 인제스트와 Cold/Warm 비교용 독립 File_IDX를 할당한다."""

        file_idx = self._allocate_file_idx()
        return source.clone(
            case_id=f"{purpose}-{source.file_type}-{source.profile_name}-{ordinal:04d}",
            file_idx=file_idx,
            group=purpose,
        )

    def _generate_one(
        self,
        *,
        group: str,
        file_type: SupportedFormat,
        profile: FixtureProfile,
    ) -> GeneratedFixture:
        file_idx = self._allocate_file_idx()
        case_id = f"{group}-{file_type}-{profile.name}"
        extension = _EXTENSION_BY_FORMAT[file_type]
        output_path = self._output_directory / f"{file_idx}-{case_id}{extension}"

        answer_fact = (
            f"Document {file_idx} uses format {file_type.upper()} and profile {profile.name}. "
            f"It contains {profile.text_units} declared text units and "
            f"{profile.image_count} declared OCR images."
        )
        text_anchor = f"JIPSA-PERFORMANCE-{file_idx}-{file_type.upper()}-{profile.name.upper()}"
        ocr_anchor = f"JIPSA PERFORMANCE OCR FILE {file_idx}"
        search_query = ocr_anchor if profile.ocr_only else text_anchor

        text_units = tuple(
            self._build_text_unit(
                case_id=case_id,
                file_idx=file_idx,
                file_type=file_type,
                profile=profile,
                unit_index=unit_index,
                answer_fact=answer_fact,
                text_anchor=text_anchor,
            )
            for unit_index in range(profile.text_units)
        )
        image_paths = tuple(
            self._create_ocr_image(
                case_id=case_id,
                file_idx=file_idx,
                image_index=image_index,
                image_count=profile.image_count,
            )
            for image_index in range(profile.image_count)
        )

        generators = {
            "pdf": self._write_pdf,
            "docx": self._write_docx,
            "pptx": self._write_pptx,
            "xlsx": self._write_xlsx,
            "txt": self._write_txt,
        }
        generators[file_type](
            output_path=output_path,
            text_units=text_units,
            image_paths=image_paths,
            title=case_id,
        )

        if not output_path.is_file() or output_path.stat().st_size <= 0:
            raise RuntimeError(f"Performance fixture was not created: {output_path}")

        return GeneratedFixture(
            case_id=case_id,
            group=group,
            file_idx=file_idx,
            file_type=file_type,
            profile_name=profile.name,
            content_origin=profile.content_origin,
            path=output_path,
            content_type=_CONTENT_TYPE_BY_FORMAT[file_type],
            text_units=profile.text_units,
            image_count=profile.image_count,
            search_query=search_query,
            answer_fact=answer_fact,
        )

    def _allocate_file_idx(self) -> int:
        file_idx = self._next_file_idx
        self._next_file_idx += 1
        return file_idx

    @staticmethod
    def _build_text_unit(
        *,
        case_id: str,
        file_idx: int,
        file_type: SupportedFormat,
        profile: FixtureProfile,
        unit_index: int,
        answer_fact: str,
        text_anchor: str,
    ) -> str:
        """압축으로 크기 차이가 사라지지 않도록 고유 토큰을 포함한 텍스트를 만든다."""

        repetitions: list[str] = []
        for repetition_index in range(profile.repetitions_per_unit):
            digest = hashlib.sha256(
                f"{case_id}:{unit_index}:{repetition_index}".encode()
            ).hexdigest()
            repetitions.append(
                " ".join(
                    (
                        f"UNIT-{unit_index:04d}",
                        f"SEGMENT-{repetition_index:04d}",
                        f"FORMAT-{file_type.upper()}",
                        f"PROFILE-{profile.name.upper()}",
                        f"TOKEN-{digest}",
                    )
                )
            )

        prefix = (
            f"{text_anchor}. {answer_fact} "
            if unit_index == 0
            else f"JIPSA performance continuation for file {file_idx}. "
        )
        return prefix + " ".join(repetitions)

    def _create_ocr_image(
        self,
        *,
        case_id: str,
        file_idx: int,
        image_index: int,
        image_count: int,
    ) -> Path:
        """EasyOCR이 읽기 쉬운 고대비 이미지와 결정적 배경 패턴을 생성한다."""

        width = 1600
        height = 900
        image = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(image)

        # 문서 크기가 이미지 수에 따라 실제로 증가하도록 텍스트 영역 밖에 결정적인
        # 회색 패턴을 배치한다. OCR 핵심 문구는 흰 배경 중앙에 유지한다.
        seed = hashlib.sha256(f"{case_id}:{image_index}".encode()).digest()
        for block_index in range(180):
            source = seed[block_index % len(seed)]
            block_x = 20 + ((source * 37 + block_index * 53) % 1540)
            block_y = 520 + ((source * 19 + block_index * 29) % 340)
            size = 3 + (source % 10)
            shade = 205 + (source % 40)
            draw.rectangle(
                (block_x, block_y, block_x + size, block_y + size),
                fill=(shade, shade, shade),
            )

        font = _load_ocr_font(size=64)
        small_font = _load_ocr_font(size=42)
        lines = (
            "JIPSA PERFORMANCE OCR",
            f"FILE {file_idx}",
            f"IMAGE {image_index + 1} OF {image_count}",
            "RESOURCE LIMIT BENCHMARK",
        )

        text_y = 80.0
        for line_index, line in enumerate(lines):
            selected_font = font if line_index == 0 else small_font
            bounds = draw.textbbox((0, 0), line, font=selected_font)
            text_width = float(bounds[2] - bounds[0])
            text_x = max((float(width) - text_width) / 2.0, 40.0)
            draw.text((text_x, text_y), line, fill="black", font=selected_font)
            text_y += 105.0

        output_path = self._image_directory / f"{file_idx}-{image_index + 1:03d}.png"
        image.save(output_path, format="PNG", optimize=False)
        return output_path

    @staticmethod
    def _write_pdf(
        *,
        output_path: Path,
        text_units: tuple[str, ...],
        image_paths: tuple[Path, ...],
        title: str,
    ) -> None:
        document = fitz.open()
        try:
            if text_units:
                units_per_page = 1
                for page_start in range(0, len(text_units), units_per_page):
                    page = document.new_page(width=595, height=842)
                    page.insert_text((40, 50), title, fontsize=14)
                    page_text = "\n\n".join(
                        textwrap.fill(unit, width=92)
                        for unit in text_units[page_start : page_start + units_per_page]
                    )
                    page.insert_textbox(
                        fitz.Rect(40, 75, 555, 810),
                        page_text,
                        fontsize=8,
                        lineheight=1.15,
                    )

            for image_path in image_paths:
                page = document.new_page(width=595, height=842)
                page.insert_image(
                    fitz.Rect(35, 120, 560, 680),
                    filename=str(image_path),
                    keep_proportion=True,
                )

            if document.page_count == 0:
                document.new_page(width=595, height=842)
            document.set_metadata({"title": title})
            document.save(output_path, garbage=4, deflate=True)
        finally:
            document.close()

    @staticmethod
    def _write_docx(
        *,
        output_path: Path,
        text_units: tuple[str, ...],
        image_paths: tuple[Path, ...],
        title: str,
    ) -> None:
        document = Document()
        if text_units:
            document.add_heading(title, level=1)
            for unit_index, unit in enumerate(text_units):
                document.add_heading(f"Unit {unit_index + 1}", level=2)
                document.add_paragraph(unit)
        for image_index, image_path in enumerate(image_paths):
            document.add_picture(str(image_path), width=DocxInches(6.2))
            # 캡션은 OCR-only 비교에서 일반 텍스트 근거가 되지 않도록 추가하지 않는다.
            if image_index < len(image_paths) - 1:
                # python-docx의 타입 Stub이 반환 타입을 제공하지 않으므로 이 호출 경계만
                # 제한적으로 무시한다. 실제 동작과 생성 문서 구조는 단위 테스트로 검증한다.
                document.add_page_break()  # type: ignore[no-untyped-call]
        document.core_properties.title = title
        document.save(str(output_path))

    @staticmethod
    def _write_pptx(
        *,
        output_path: Path,
        text_units: tuple[str, ...],
        image_paths: tuple[Path, ...],
        title: str,
    ) -> None:
        presentation = Presentation()
        # 기본 생성되는 첫 슬라이드는 없으므로 text/image 조건에 맞춰 명시적으로 추가한다.
        blank_layout = presentation.slide_layouts[6]

        for unit_index, unit in enumerate(text_units):
            slide = presentation.slides.add_slide(blank_layout)
            title_box = slide.shapes.add_textbox(
                PptxInches(0.5),
                PptxInches(0.3),
                PptxInches(9.0),
                PptxInches(0.6),
            )
            title_box.text_frame.text = f"{title} - Unit {unit_index + 1}"
            body_box = slide.shapes.add_textbox(
                PptxInches(0.5),
                PptxInches(1.0),
                PptxInches(9.0),
                PptxInches(6.0),
            )
            body_box.text_frame.word_wrap = True
            body_box.text_frame.text = unit

        for image_path in image_paths:
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                PptxInches(0.7),
                PptxInches(0.8),
                width=PptxInches(8.6),
            )

        if len(presentation.slides) == 0:
            presentation.slides.add_slide(blank_layout)
        presentation.core_properties.title = title
        presentation.save(str(output_path))

    @staticmethod
    def _write_xlsx(
        *,
        output_path: Path,
        text_units: tuple[str, ...],
        image_paths: tuple[Path, ...],
        title: str,
    ) -> None:
        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Performance"

        if text_units:
            worksheet["A1"] = title
            worksheet["A2"] = "Unit"
            worksheet["B2"] = "Content"
            for unit_index, unit in enumerate(text_units, start=1):
                worksheet.cell(row=unit_index + 2, column=1, value=unit_index)
                worksheet.cell(row=unit_index + 2, column=2, value=unit)
            worksheet.column_dimensions["A"].width = 12
            worksheet.column_dimensions["B"].width = 120

        for image_index, image_path in enumerate(image_paths):
            image = OpenpyxlImage(str(image_path))
            image.width = 800
            image.height = 450
            start_row = 1 + (image_index * 28)
            worksheet.add_image(image, f"D{start_row}")

        workbook.properties.title = title
        workbook.save(output_path)
        workbook.close()

    @staticmethod
    def _write_txt(
        *,
        output_path: Path,
        text_units: tuple[str, ...],
        image_paths: tuple[Path, ...],
        title: str,
    ) -> None:
        if image_paths:
            raise ValueError("TXT performance fixtures do not support embedded images.")
        payload = "\n\n".join((title, *text_units))
        output_path.write_text(payload, encoding="utf-8", newline="\n")


def _load_ocr_font(*, size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """Windows와 Linux에서 사용 가능한 명확한 TrueType Font를 순서대로 찾는다."""

    candidates: Iterable[Path] = (
        Path("C:/Windows/Fonts/arialbd.ttf"),
        Path("C:/Windows/Fonts/malgunbd.ttf"),
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
    )
    for path in candidates:
        if path.is_file():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()
