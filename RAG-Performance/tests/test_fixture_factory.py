"""합성 Fixture가 다섯 문서 형식과 OCR-only 조건을 생성하는지 검증한다."""

from pathlib import Path

import fitz  # type: ignore[import-untyped]
from docx import Document
from openpyxl import load_workbook  # type: ignore[import-untyped]
from pptx import Presentation

from jipsa_rag_benchmark.fixture_factory import FixtureFactory
from jipsa_rag_benchmark.models import (
    AnswerLoadPlan,
    BenchmarkPlan,
    FixtureMatrixEntry,
    FixtureProfile,
    LoadPlan,
    SaturationPolicy,
    SearchPlan,
)


def test_fixture_factory_generates_five_text_formats_and_ocr_pdf(tmp_path: Path) -> None:
    fixtures = FixtureFactory(
        output_directory=tmp_path / "fixtures",
        file_idx_start=1599000,
    ).generate(_small_test_plan())

    assert len(fixtures) == 6
    assert {fixture.file_type for fixture in fixtures} == {
        "pdf",
        "docx",
        "pptx",
        "xlsx",
        "txt",
    }
    assert {fixture.content_origin for fixture in fixtures} == {"text", "ocr"}

    text_fixtures = {
        fixture.file_type: fixture
        for fixture in fixtures
        if fixture.content_origin == "text"
    }
    assert text_fixtures["pdf"].search_query in _pdf_text(text_fixtures["pdf"].path)
    assert text_fixtures["docx"].search_query in _docx_text(text_fixtures["docx"].path)
    assert text_fixtures["pptx"].search_query in _pptx_text(text_fixtures["pptx"].path)
    assert text_fixtures["xlsx"].search_query in _xlsx_text(text_fixtures["xlsx"].path)
    assert text_fixtures["txt"].search_query in text_fixtures["txt"].path.read_text(
        encoding="utf-8"
    )

    ocr_pdf = next(
        fixture
        for fixture in fixtures
        if fixture.file_type == "pdf" and fixture.content_origin == "ocr"
    )
    with fitz.open(ocr_pdf.path) as document:
        assert document.page_count == 2
        assert all(page.get_text().strip() == "" for page in document)
        assert sum(len(page.get_images(full=True)) for page in document) == 2


def _small_test_plan() -> BenchmarkPlan:
    load = LoadPlan(concurrency_levels=(1,), requests_per_level=1)
    return BenchmarkPlan(
        schema_version=1,
        benchmark_name="standalone-fixture-unit-test",
        test_user_idx=159900,
        file_idx_start=1599000,
        sample_interval_seconds=1.0,
        docker_sample_interval_seconds=2.0,
        request_timeout_seconds=30.0,
        warmup_requests=0,
        fixture_profiles=(
            FixtureProfile(
                name="tiny_text",
                text_units=2,
                repetitions_per_unit=2,
                image_count=0,
                ocr_only=False,
            ),
            FixtureProfile(
                name="tiny_ocr",
                text_units=0,
                repetitions_per_unit=0,
                image_count=2,
                ocr_only=True,
            ),
        ),
        fixture_matrix=(
            FixtureMatrixEntry(
                group="format-coverage",
                formats=("pdf", "docx", "pptx", "xlsx", "txt"),
                profiles=("tiny_text",),
            ),
            FixtureMatrixEntry(
                group="ocr-comparison",
                formats=("pdf",),
                profiles=("tiny_ocr",),
            ),
        ),
        ingest=load,
        search=SearchPlan(top_k=5, score_threshold=None, load=load),
        answers=AnswerLoadPlan(enabled=False, lookup=load, synthesis=load),
        saturation=SaturationPolicy(
            max_error_rate=0.01,
            throughput_gain_floor_percent=5.0,
            p95_growth_trigger_percent=20.0,
        ),
    )


def _pdf_text(path: Path) -> str:
    with fitz.open(path) as document:
        return "\n".join(page.get_text() for page in document)


def _docx_text(path: Path) -> str:
    document = Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _pptx_text(path: Path) -> str:
    presentation = Presentation(path)
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def _xlsx_text(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        return "\n".join(
            str(cell.value)
            for worksheet in workbook.worksheets
            for row in worksheet.iter_rows()
            for cell in row
            if cell.value is not None
        )
    finally:
        workbook.close()
