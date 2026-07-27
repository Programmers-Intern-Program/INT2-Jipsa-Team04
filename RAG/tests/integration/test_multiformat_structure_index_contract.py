"""다섯 문서 형식의 파싱·구조 청킹·저장 payload 통합 계약을 테스트한다.

외부 MySQL, Qdrant 또는 TEI 서버를 사용하지 않고 실제 형식 파서가 만든 위치
메타데이터가 구조화 청킹과 두 저장소 변환 경계를 통과해 보존되는지 검증한다.
"""

import hashlib
import json
from collections.abc import Awaitable, Callable, Mapping
from pathlib import Path
from typing import cast

import pytest
from docx import Document
from openpyxl import Workbook  # type: ignore[import-untyped]
from pptx import Presentation
from pptx.util import Inches

from jipsa_rag.infrastructure.chunking.models import ChunkingContext
from jipsa_rag.infrastructure.chunking.structured import StructuredDocumentChunker
from jipsa_rag.infrastructure.document.models import (
    DocumentType,
    ParsedDocument,
    ParsedDocumentUnit,
)
from jipsa_rag.infrastructure.document.parsers.docx import DocxDocumentParser
from jipsa_rag.infrastructure.document.parsers.pptx import PptxDocumentParser
from jipsa_rag.infrastructure.document.parsers.txt import TxtDocumentParser
from jipsa_rag.infrastructure.document.parsers.xlsx import XlsxDocumentParser
from jipsa_rag.infrastructure.embedding.models import EmbeddedChunk, EmbeddedDocument
from jipsa_rag.infrastructure.indexing.local_repository import _build_chunk_parameters
from jipsa_rag.infrastructure.indexing.models import DocumentIndexMetadata
from jipsa_rag.infrastructure.indexing.qdrant_store import _build_payload

ParsedFactory = Callable[[Path], Awaitable[tuple[ParsedDocument, str, str]]]


async def _pdf_document(_: Path) -> tuple[ParsedDocument, str, str]:
    """기존 PDF 페이지 메타데이터 계약을 사용하는 메모리 문서를 생성한다."""

    return (
        ParsedDocument(
            file_type=DocumentType.PDF,
            units=(
                ParsedDocumentUnit(
                    text="PDF 첫 페이지 본문",
                    source_metadata={"page_number": 1},
                ),
            ),
            document_metadata={"page_count": 1},
        ),
        "PDF_TEXT",
        "1.0.0",
    )


async def _docx_document(tmp_path: Path) -> tuple[ParsedDocument, str, str]:
    """제목과 문단이 포함된 실제 DOCX를 생성하고 파싱한다."""

    file_path = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("운영 가이드", level=1)
    document.add_paragraph("서비스를 시작하고 상태를 확인한다.")
    document.save(str(file_path))
    parser = DocxDocumentParser()
    return await parser.parse(file_path), parser.parser_type, parser.parser_version


async def _pptx_document(tmp_path: Path) -> tuple[ParsedDocument, str, str]:
    """위치가 지정된 텍스트 도형이 포함된 실제 PPTX를 생성하고 파싱한다."""

    file_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.text = "PPTX 구조 메타데이터"
    presentation.save(str(file_path))
    parser = PptxDocumentParser()
    return await parser.parse(file_path), parser.parser_type, parser.parser_version


async def _xlsx_document(tmp_path: Path) -> tuple[ParsedDocument, str, str]:
    """시트와 셀 범위가 있는 실제 XLSX를 생성하고 파싱한다."""

    file_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "요약"
    sheet.append(["항목", "값"])
    sheet.append(["진행률", 80])
    workbook.save(file_path)
    workbook.close()
    parser = XlsxDocumentParser()
    return await parser.parse(file_path), parser.parser_type, parser.parser_version


async def _txt_document(tmp_path: Path) -> tuple[ParsedDocument, str, str]:
    """여러 줄 UTF-8 TXT를 생성하고 파싱한다."""

    file_path = tmp_path / "sample.txt"
    file_path.write_text("첫 번째 줄\n두 번째 줄", encoding="utf-8")
    parser = TxtDocumentParser()
    return await parser.parse(file_path), parser.parser_type, parser.parser_version


@pytest.mark.asyncio
@pytest.mark.parametrize(
    ("factory", "expected_location_key"),
    (
        (_pdf_document, "page_number"),
        (_docx_document, "section_index"),
        (_pptx_document, "shape_path"),
        (_xlsx_document, "cell_range"),
        (_txt_document, "line_number"),
    ),
)
async def test_position_metadata_survives_parse_chunk_db_and_qdrant_contract(
    tmp_path: Path,
    factory: ParsedFactory,
    expected_location_key: str,
) -> None:
    """형식별 위치가 Local DB JSON과 Qdrant payload까지 손실 없이 전달된다."""

    parsed_document, parser_type, parser_version = await factory(tmp_path)
    file_bytes = parsed_document.text.encode("utf-8")
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    chunked_document = await StructuredDocumentChunker().chunk(
        document=parsed_document,
        context=ChunkingContext(
            users_idx=1,
            file_idx=2,
            file_hash=file_hash,
            parser_version=parser_version,
            embedding_model="test/embedding-model",
            index_version=2,
        ),
    )
    first_chunk = chunked_document.chunks[0]
    embedded_document = EmbeddedDocument(
        embedding_model="test/embedding-model",
        embedding_dim=3,
        chunks=(EmbeddedChunk(chunk=first_chunk, embedding=(0.1, 0.2, 0.3)),),
    )
    metadata = DocumentIndexMetadata(
        users_idx=1,
        file_idx=2,
        folder_idx=None,
        file_name=f"sample.{parsed_document.file_type.value.lower()}",
        file_type=parsed_document.file_type,
        file_hash=file_hash,
        index_version=2,
        parser_type=parser_type,
        parser_version=parser_version,
    )

    db_parameters = _build_chunk_parameters(
        rag_document_idx=3,
        metadata=metadata,
        embedding_model=embedded_document.embedding_model,
        embedded_chunk=embedded_document.chunks[0],
    )
    qdrant_payload = _build_payload(
        rag_document_idx=3,
        metadata=metadata,
        embedded_document=embedded_document,
        embedded_chunk_index=0,
        chunk=first_chunk,
        is_active=True,
        created_at="2026-07-27T00:00:00+00:00",
    )

    qdrant_source_metadata = cast(
        Mapping[str, object],
        qdrant_payload["source_metadata"],
    )
    db_source_metadata = cast(
        dict[str, object],
        json.loads(cast(str, db_parameters["source_metadata"])),
    )

    assert expected_location_key in first_chunk.source_metadata
    assert expected_location_key in qdrant_source_metadata
    assert expected_location_key in db_source_metadata
    assert qdrant_payload["chunk_id"] == db_parameters["chunk_id"]
    assert qdrant_payload["content_hash"] == db_parameters["content_hash"]
    assert qdrant_payload["chunking_strategy"] == "STRUCTURED_DOCUMENT"
