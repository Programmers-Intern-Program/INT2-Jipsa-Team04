"""실제 PowerPoint·Excel COM에서 다운로드 임시 확장자 호환성을 검증한다."""

from __future__ import annotations

import os
import sys
from pathlib import Path
from typing import Any, cast

import pytest
from openpyxl import Workbook  # type: ignore[import-untyped]
from openpyxl.chart import BarChart, Reference  # type: ignore[import-untyped]
from openpyxl.drawing.image import Image as XlsxImage  # type: ignore[import-untyped]
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from jipsa_rag.core.document_processing import DocumentProcessingSettings
from jipsa_rag.infrastructure.document.images.download_safe_pptx import (
    DownloadSafePptxImageExtractor,
)
from jipsa_rag.infrastructure.document.images.download_safe_xlsx import (
    DownloadSafeXlsxImageExtractor,
)
from jipsa_rag.infrastructure.document.images.models import DocumentImageKind
from jipsa_rag.infrastructure.document.rendering import MicrosoftOfficeRenderClient

_OFFICE_COM_INTEGRATION_ENV = "JIPSA_RAG_RUN_OFFICE_COM_INTEGRATION"


def _require_office_com_integration() -> None:
    """명시적으로 활성화된 Windows 사용자 세션에서만 실제 Office COM을 실행한다."""

    if sys.platform != "win32":
        pytest.skip("Microsoft Office COM integration requires Windows.")
    if os.environ.get(_OFFICE_COM_INTEGRATION_ENV) != "1":
        pytest.skip(f"Set {_OFFICE_COM_INTEGRATION_ENV}=1 to run Office COM tests.")


def _settings() -> DocumentProcessingSettings:
    """실제 Office 렌더링 통합 테스트에 필요한 제한값을 생성한다."""

    return DocumentProcessingSettings(
        office_rendering_enabled=True,
        office_rendering_provider="microsoft_office_com",
        office_com_require_interactive_session=True,
        office_render_max_concurrency=1,
        office_render_timeout_seconds=60,
        image_decorative_filter_enabled=False,
        image_max_count_per_document=50,
        image_max_bytes=50 * 1024 * 1024,
        image_max_pixels=80_000_000,
        ocr_enabled=False,
        ocr_gpu=False,
        ocr_gpu_required=False,
        _env_file=None,
    )


def _create_png(path: Path) -> None:
    """PPTX와 XLSX에 공통으로 삽입할 검색 가능한 크기의 PNG를 생성한다."""

    Image.new("RGB", (600, 240), "white").save(path, format="PNG")


def _create_pptx(path: Path, image_path: Path) -> None:
    """삽입 그림과 실제 PowerPoint 차트를 포함하는 PPTX를 생성한다."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        Inches(0.5),
        Inches(0.5),
        width=Inches(3),
    )

    chart_data_factory = cast(Any, ChartData)
    chart_data = chart_data_factory()
    chart_data.categories = ["A", "B", "C"]
    cast(Any, chart_data).add_series("Series", (1, 3, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(4),
        Inches(1),
        Inches(5),
        Inches(4),
        chart_data,
    )
    presentation.save(str(path))


def _create_xlsx(path: Path, image_path: Path) -> None:
    """삽입 그림과 실제 Excel 차트를 포함하는 XLSX를 생성한다."""

    workbook = Workbook()
    worksheet = workbook.active
    if worksheet is None:
        raise AssertionError("A new workbook must contain an active worksheet.")

    worksheet.title = "Sheet1"
    worksheet.append(["Category", "Value"])
    worksheet.append(["A", 1])
    worksheet.append(["B", 3])
    worksheet.append(["C", 2])
    worksheet.add_image(XlsxImage(str(image_path)), "D2")

    chart = BarChart()
    chart.add_data(
        Reference(worksheet, min_col=2, min_row=1, max_row=4),
        titles_from_data=True,
    )
    chart.set_categories(
        Reference(worksheet, min_col=1, min_row=2, max_row=4),
    )
    worksheet.add_chart(chart, "D15")
    workbook.save(str(path))
    workbook.close()


@pytest.mark.asyncio
async def test_document_extension_renders_pptx_and_xlsx_visuals_with_office_com(
    tmp_path: Path,
) -> None:
    """실제 ``*.document`` 입력에서도 PowerPoint·Excel 시각 요소를 모두 렌더링한다."""

    _require_office_com_integration()

    image_path = tmp_path / "sample.png"
    pptx_path = tmp_path / "sample.pptx"
    xlsx_path = tmp_path / "sample.xlsx"
    downloaded_pptx_path = tmp_path / "downloaded-pptx.document"
    downloaded_xlsx_path = tmp_path / "downloaded-xlsx.document"
    _create_png(image_path)
    _create_pptx(pptx_path, image_path)
    _create_xlsx(xlsx_path, image_path)

    # 운영 HttpFileDownloader와 동일하게 정상 Office 바이트를 확장자가 제거된
    # 서로 다른 임시 파일에 저장한다.
    downloaded_pptx_path.write_bytes(pptx_path.read_bytes())
    downloaded_xlsx_path.write_bytes(xlsx_path.read_bytes())

    settings = _settings()
    renderer = MicrosoftOfficeRenderClient(settings)
    pptx_result = await DownloadSafePptxImageExtractor(settings, renderer).extract(
        downloaded_pptx_path
    )
    xlsx_result = await DownloadSafeXlsxImageExtractor(settings, renderer).extract(
        downloaded_xlsx_path
    )

    assert pptx_result.document_metadata["office_renderer_available"] is True
    assert xlsx_result.document_metadata["office_renderer_available"] is True

    assert any(image.kind is DocumentImageKind.PPTX_PICTURE for image in pptx_result.images)
    assert any(image.kind is DocumentImageKind.PPTX_CHART_RENDER for image in pptx_result.images)
    assert any(image.kind is DocumentImageKind.XLSX_PICTURE for image in xlsx_result.images)
    assert any(image.kind is DocumentImageKind.XLSX_CHART_RENDER for image in xlsx_result.images)

    # 안전 어댑터는 실제 인제스트 입력 파일을 삭제하거나 수정하지 않는다.
    assert downloaded_pptx_path.read_bytes() == pptx_path.read_bytes()
    assert downloaded_xlsx_path.read_bytes() == xlsx_path.read_bytes()
