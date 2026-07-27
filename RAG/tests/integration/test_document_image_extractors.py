"""실제 문서 파일을 생성하여 형식별 이미지 추출 경로를 검증한다."""

from __future__ import annotations

import os
import sys
import zipfile
from pathlib import Path
from typing import Any, cast

# PyMuPDF는 py.typed marker와 공식 type stub을 제공하지 않는다. 테스트는 실제
# PDF 파일 생성과 추출 결과로 런타임 계약을 검증하므로 이 import 경계만 격리한다.
import fitz  # type: ignore[import-untyped]
import pytest
from docx import Document

# openpyxl은 py.typed를 제공하지 않고, 외부 types-openpyxl 패키지는 실제
# Worksheet와 Chartsheet 호출 계약을 부정확하게 선언한다. 테스트 fixture 생성에
# 사용하는 외부 라이브러리 import 경계만 격리하고 실제 결과는 통합 테스트로 검증한다.
from openpyxl import Workbook  # type: ignore[import-untyped]
from openpyxl.chart import BarChart, Reference  # type: ignore[import-untyped]
from openpyxl.drawing.image import Image as XlsxImage  # type: ignore[import-untyped]
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from jipsa_rag.core.document_processing import DocumentProcessingSettings
from jipsa_rag.infrastructure.document.images.docx import DocxImageExtractor
from jipsa_rag.infrastructure.document.images.models import DocumentImageKind
from jipsa_rag.infrastructure.document.images.pdf import PdfImageExtractor
from jipsa_rag.infrastructure.document.images.pptx import PptxImageExtractor
from jipsa_rag.infrastructure.document.images.xlsx import XlsxImageExtractor
from jipsa_rag.infrastructure.document.rendering import MicrosoftOfficeRenderClient

_OFFICE_COM_INTEGRATION_ENV = "JIPSA_RAG_RUN_OFFICE_COM_INTEGRATION"


@pytest.fixture
def sample_png(tmp_path: Path) -> Path:
    """문서에 삽입할 작은 검색용 PNG를 생성한다."""

    path = tmp_path / "sample-image.png"
    image = Image.new("RGB", (600, 240), "white")
    draw = ImageDraw.Draw(image)
    draw.text((30, 80), "OCR TEST 123", fill="black")
    image.save(path)
    return path


def _settings(*, office_rendering_enabled: bool = False) -> DocumentProcessingSettings:
    """통합 테스트에 필요한 자원 한계와 Office 렌더링 정책을 생성한다."""

    return DocumentProcessingSettings(
        office_rendering_enabled=office_rendering_enabled,
        office_rendering_provider="microsoft_office_com",
        office_com_require_interactive_session=True,
        office_render_max_concurrency=1,
        office_render_timeout_seconds=60,
        image_max_count_per_document=50,
        image_max_bytes=50 * 1024 * 1024,
        image_max_pixels=80_000_000,
        _env_file=None,
    )


def _require_office_com_integration() -> None:
    """명시적으로 활성화된 Windows 로컬 세션에서만 실제 Office COM을 실행한다."""

    if sys.platform != "win32":
        pytest.skip("Microsoft Office COM integration requires Windows.")
    if os.environ.get(_OFFICE_COM_INTEGRATION_ENV) != "1":
        pytest.skip(f"Set {_OFFICE_COM_INTEGRATION_ENV}=1 to run Office COM tests.")


@pytest.mark.asyncio
async def test_pdf_extractor_returns_embedded_and_scan_page_images(
    tmp_path: Path,
    sample_png: Path,
) -> None:
    """텍스트 페이지 이미지는 개별 추출하고 이미지 전용 페이지는 전체 렌더한다."""

    pdf_path = tmp_path / "sample.pdf"
    document = fitz.open()
    text_page = document.new_page(width=600, height=800)
    text_page.insert_text((40, 50), "Text layer page")
    text_page.insert_image(
        fitz.Rect(40, 100, 560, 308),
        filename=str(sample_png),
    )
    scan_page = document.new_page(width=600, height=800)
    scan_page.insert_image(scan_page.rect, filename=str(sample_png))
    document.save(str(pdf_path))
    document.close()

    result = await PdfImageExtractor(_settings()).extract(pdf_path)

    assert any(image.kind is DocumentImageKind.PDF_EMBEDDED for image in result.images)
    assert any(image.kind is DocumentImageKind.PDF_PAGE_RENDER for image in result.images)
    assert result.document_metadata["image_only_page_count"] == 1
    scan_render = next(
        image for image in result.images if image.kind is DocumentImageKind.PDF_PAGE_RENDER
    )
    assert scan_render.source_metadata["page_number"] == 2


@pytest.mark.asyncio
async def test_docx_extractor_distinguishes_inline_and_floating_images(
    tmp_path: Path,
    sample_png: Path,
) -> None:
    """동일 그림을 inline과 anchor XML로 각각 저장해 배치 유형을 검증한다."""

    inline_path = tmp_path / "inline.docx"
    document = Document()
    document.add_paragraph("이전 문단")
    paragraph = document.add_paragraph("이미지 문단 ")
    paragraph.add_run().add_picture(str(sample_png), width=Inches(3))
    document.add_paragraph("다음 문단")
    document.save(str(inline_path))

    floating_path = tmp_path / "floating.docx"
    _rewrite_inline_as_anchor(inline_path, floating_path)

    extractor = DocxImageExtractor(_settings())
    inline_result = await extractor.extract(inline_path)
    floating_result = await extractor.extract(floating_path)

    assert inline_result.images[0].kind is DocumentImageKind.DOCX_INLINE
    assert floating_result.images[0].kind is DocumentImageKind.DOCX_FLOATING
    assert inline_result.images[0].source_metadata["paragraph_index"] == 2
    assert inline_result.images[0].context_before == "이전 문단"
    assert inline_result.images[0].context_after == "다음 문단"


@pytest.mark.asyncio
async def test_pptx_and_xlsx_extract_pictures_and_rendered_charts(
    tmp_path: Path,
    sample_png: Path,
) -> None:
    """PowerPoint·Excel COM으로 그림과 차트 렌더 이미지를 함께 추출한다."""

    _require_office_com_integration()
    settings = _settings(office_rendering_enabled=True)
    renderer = MicrosoftOfficeRenderClient(settings)
    pptx_path = tmp_path / "sample.pptx"
    xlsx_path = tmp_path / "sample.xlsx"
    _create_pptx(pptx_path, sample_png)
    _create_xlsx(xlsx_path, sample_png)

    pptx_result = await PptxImageExtractor(settings, renderer).extract(pptx_path)
    xlsx_result = await XlsxImageExtractor(settings, renderer).extract(xlsx_path)

    assert pptx_result.document_metadata["office_renderer_available"] is True
    assert xlsx_result.document_metadata["office_renderer_available"] is True
    assert any(image.kind is DocumentImageKind.PPTX_PICTURE for image in pptx_result.images)
    assert any(image.kind is DocumentImageKind.PPTX_CHART_RENDER for image in pptx_result.images)
    assert any(image.kind is DocumentImageKind.XLSX_PICTURE for image in xlsx_result.images)
    assert any(image.kind is DocumentImageKind.XLSX_CHART_RENDER for image in xlsx_result.images)


@pytest.mark.asyncio
async def test_xlsx_chart_render_is_bound_to_actual_chart_sheet(
    tmp_path: Path,
) -> None:
    """앞 시트의 길이와 무관하게 차트를 실제 시트·anchor 메타데이터에 연결한다."""

    _require_office_com_integration()
    workbook_path = tmp_path / "multi-sheet-chart.xlsx"
    workbook = Workbook()
    first_sheet = workbook.active
    assert first_sheet is not None
    first_sheet.title = "LongData"
    first_sheet.append(["Index", "Value"])
    for value in range(1, 151):
        first_sheet.append([value, value * 2])

    chart_sheet = workbook.create_sheet("ChartTarget")
    chart_sheet.append(["Category", "Value"])
    chart_sheet.append(["A", 1])
    chart_sheet.append(["B", 3])
    chart_sheet.append(["C", 2])
    chart = BarChart()
    chart.add_data(
        Reference(chart_sheet, min_col=2, min_row=1, max_row=4),
        titles_from_data=True,
    )
    chart.set_categories(
        Reference(chart_sheet, min_col=1, min_row=2, max_row=4),
    )
    chart_sheet.add_chart(chart, "D15")
    workbook.save(str(workbook_path))
    workbook.close()

    settings = _settings(office_rendering_enabled=True)
    result = await XlsxImageExtractor(
        settings,
        MicrosoftOfficeRenderClient(settings),
    ).extract(workbook_path)

    assert result.document_metadata["office_renderer_available"] is True
    chart_image = next(
        image for image in result.images if image.kind is DocumentImageKind.XLSX_CHART_RENDER
    )
    assert chart_image.source_metadata["sheet_index"] == 2
    assert chart_image.source_metadata["sheet_name"] == "ChartTarget"
    assert chart_image.source_metadata["chart_index"] == 1
    assert chart_image.source_metadata["anchor_cell"] == "D15"
    assert chart_image.source_metadata["cell_range"] != "RENDERED_SHEET_PAGE"


def _rewrite_inline_as_anchor(source: Path, destination: Path) -> None:
    """테스트 DOCX의 첫 wp:inline 태그를 wp:anchor로 바꾸어 floating fixture를 만든다."""

    with zipfile.ZipFile(source) as input_package:
        members = {
            info.filename: input_package.read(info.filename) for info in input_package.infolist()
        }

    document_xml = members["word/document.xml"]
    document_xml = document_xml.replace(b"<wp:inline ", b"<wp:anchor ", 1)
    document_xml = document_xml.replace(b"</wp:inline>", b"</wp:anchor>", 1)
    members["word/document.xml"] = document_xml

    with zipfile.ZipFile(destination, "w", compression=zipfile.ZIP_DEFLATED) as output:
        for name, content in members.items():
            output.writestr(name, content)


def _create_pptx(path: Path, sample_png: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(sample_png),
        Inches(0.5),
        Inches(0.5),
        width=Inches(3),
    )
    chart_data_factory = cast(Any, ChartData)
    chart_data = chart_data_factory()
    chart_data.categories = ["A", "B", "C"]
    cast(Any, chart_data).add_series("Series", (1, 3, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(4),
        Inches(1),
        Inches(5),
        Inches(4),
        chart_data,
    )
    presentation.save(str(path))


def _create_xlsx(path: Path, sample_png: Path) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    assert worksheet is not None
    worksheet.title = "Sheet1"
    worksheet.append(["Category", "Value"])
    worksheet.append(["A", 1])
    worksheet.append(["B", 3])
    worksheet.append(["C", 2])
    worksheet.add_image(XlsxImage(str(sample_png)), "D2")

    chart = BarChart()
    chart.add_data(
        Reference(worksheet, min_col=2, min_row=1, max_row=4),
        titles_from_data=True,
    )
    chart.set_categories(
        Reference(worksheet, min_col=1, min_row=2, max_row=4),
    )
    worksheet.add_chart(chart, "D15")
    workbook.save(str(path))
    workbook.close()
