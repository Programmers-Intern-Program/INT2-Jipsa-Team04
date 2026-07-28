"""실제 DOCX·PPTX·XLSX·TXT부터 Claude 답변까지 전체 RAG 경로를 검증한다.

기존 ``test_real_pdf_rag_e2e.py``는 PDF 전용 회귀 검증으로 그대로 유지한다.
이 모듈은 나머지 네 문서 형식을 실제 파일 바이트로 생성하고 다음 경로를
운영 구현 그대로 통과시킨다.

- Backend manifest 조회
- Presigned URL 역할의 스트리밍 다운로드
- 형식별 문서 파싱
- 원본 구조를 보존하는 청킹
- CUDA TEI 임베딩
- Local RAG DB 저장
- Qdrant 벡터 및 payload 저장
- 사용자·참조문서 범위가 적용된 청크 검색
- 실제 Claude 기반 혼합 문서 답변과 SOURCE-N 출처 반환

AWS Backend와 S3의 HTTP 경계만 결정적인 ``MockTransport``로 교체한다.
Local RAG DB, Qdrant, CUDA TEI와 Claude는 실제 구성요소를 사용하므로
``JIPSA_RAG_RUN_E2E=1``을 명시한 경우에만 실행한다.
"""

from __future__ import annotations

import asyncio
import hashlib
import json
import os
import re
from collections.abc import Iterator, Mapping, Sequence
from dataclasses import dataclass, field
from io import BytesIO
from math import ceil
from pathlib import Path
from typing import Final, cast
from urllib.parse import urlsplit

import httpx2
import pytest
from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook, load_workbook  # type: ignore[import-untyped]
from pptx import Presentation
from pptx.util import Inches
from pydantic import ValidationError
from qdrant_client import AsyncQdrantClient, models
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncEngine, create_async_engine

from jipsa_rag.api.ingest import get_application_server_ingest_client
from jipsa_rag.api.v1.endpoints.file_processing import get_file_downloader
from jipsa_rag.core.config import Settings, get_settings
from jipsa_rag.core.generation_config import get_generation_settings
from jipsa_rag.infrastructure.app_server.ingest_client import (
    ApplicationServerIngestClient,
)
from jipsa_rag.infrastructure.file.downloader import HttpFileDownloader
from jipsa_rag.main import app

# ============================================================
# 실제 E2E 실행 제어 및 테스트 전용 식별자
# ============================================================

# 실제 Claude 호출과 Local 인프라 데이터 변경을 동반하므로 일반 테스트에서는
# 실행하지 않는다. 기존 PDF E2E와 동일한 Opt-in 환경 변수를 공유한다.
_RUN_ENV: Final[str] = "JIPSA_RAG_RUN_E2E"

# 기존 PDF E2E의 94_001 사용자 범위와 충돌하지 않도록 별도 사용자·파일 범위를
# 사용한다. 테스트 시작 전과 종료 후 이 범위만 선택적으로 정리한다.
_TEST_USER_IDX: Final[int] = 94_101
_TEST_FOLDER_IDX: Final[int] = 94_110
_DOCX_FILE_IDX: Final[int] = 941_001
_PPTX_FILE_IDX: Final[int] = 941_002
_XLSX_FILE_IDX: Final[int] = 941_003
_TXT_FILE_IDX: Final[int] = 941_004

_FILE_IDXS: Final[tuple[int, int, int, int]] = (
    _DOCX_FILE_IDX,
    _PPTX_FILE_IDX,
    _XLSX_FILE_IDX,
    _TXT_FILE_IDX,
)
_FILE_IDX_MIN: Final[int] = min(_FILE_IDXS)
_FILE_IDX_MAX: Final[int] = max(_FILE_IDXS)

# 현재 파일 처리 엔드포인트와 Local RAG 저장 계약의 색인 버전이다.
_INDEX_VERSION: Final[int] = 2

# DOCX·PPTX·XLSX는 운영 DocumentParserFactory에서 텍스트 파서와 이미지 추출·OCR
# 보강을 결합한 Hybrid Parser로 구성된다. 테스트 Fixture에 이미지가 없더라도
# 선택된 Parser 구현의 식별자와 버전은 HYBRID_OCR/2.0.0 계약을 유지한다.
_OCR_PARSER_VERSION: Final[str] = "2.0.0"

# 문서 및 청크 SHA-256 값의 저장 형식을 검증한다.
_SHA256_PATTERN: Final[re.Pattern[str]] = re.compile(r"^[0-9a-f]{64}$")
_SOURCE_PATTERN: Final[re.Pattern[str]] = re.compile(r"\[(SOURCE-[1-9][0-9]*)\]")

# Mock Backend가 허용할 내부 API 경로를 두 종류로 제한한다.
_BACKEND_PATH_PATTERN: Final[re.Pattern[str]] = re.compile(
    r"^/internal/files/(?P<file_idx>[1-9][0-9]*)/"
    r"(?P<operation>manifest|ingest-complete)$"
)

# 다운로드 경로의 File_IDX와 확장자를 동시에 확인한다. OOXML 형식은 모두 ZIP
# Magic Byte를 사용하므로 URL 확장자와 manifest file_type까지 일치해야 한다.
_DOWNLOAD_PATH_PATTERN: Final[re.Pattern[str]] = re.compile(
    r"^/files/(?P<file_idx>[1-9][0-9]*)\."
    r"(?P<extension>docx|pptx|xlsx|txt)$"
)


# ============================================================
# 실제 비 PDF 문서 Fixture 생성
# ============================================================


def _build_docx() -> bytes:
    """제목, 일반 문단과 표를 포함하는 실제 DOCX 바이트를 생성한다."""

    document = Document()
    document.add_heading("JIPSA DOCX E2E", level=1)
    document.add_paragraph(
        "The DOCX exact verification code is DOCX-DELTA-52. The owner is Document Operations."
    )

    # 표 추출 경로도 같은 실제 파일 안에서 실행되도록 짧은 표를 추가한다.
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Field"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Review window"
    table.cell(1, 1).text = "52 minutes"

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _build_pptx() -> bytes:
    """두 개의 텍스트 도형과 표를 포함하는 실제 PPTX 바이트를 생성한다."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    title = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(0.6),
        Inches(8.0),
        Inches(0.7),
    )
    title.text = "JIPSA PPTX E2E"

    body = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.6),
        Inches(8.0),
        Inches(1.2),
    )
    body.text = (
        "The PPTX exact verification code is PPTX-ECHO-63. The owner is Presentation Operations."
    )

    # GraphicFrame 표 추출과 도형 위치 payload가 함께 검증되도록 표를 추가한다.
    table_shape = slide.shapes.add_table(
        2,
        2,
        Inches(0.8),
        Inches(3.2),
        Inches(5.0),
        Inches(1.4),
    )
    table = table_shape.table
    table.cell(0, 0).text = "Field"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Review window"
    table.cell(1, 1).text = "63 minutes"

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_xlsx() -> bytes:
    """시트, 일반 셀, 병합 셀과 표 형태 행을 포함하는 실제 XLSX를 만든다."""

    workbook = Workbook()
    worksheet = workbook.active

    if worksheet is None:
        raise AssertionError("A new XLSX workbook must contain an active worksheet.")

    worksheet.title = "E2E"
    worksheet.merge_cells("A1:B1")
    worksheet["A1"] = "JIPSA XLSX E2E"
    worksheet["A2"] = "Exact verification code"
    worksheet["B2"] = "XLSX-FOXTROT-74"
    worksheet["A3"] = "Owner"
    worksheet["B3"] = "Spreadsheet Operations"
    worksheet["A4"] = "Review window"
    worksheet["B4"] = "74 minutes"

    buffer = BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def _build_txt() -> bytes:
    """여러 줄과 마지막 개행을 포함하는 실제 UTF-8 TXT 바이트를 생성한다."""

    text_payload = (
        "JIPSA TXT E2E\n"
        "The TXT exact verification code is TXT-GOLF-85.\n"
        "The owner is Text Operations.\n"
        "The review window is 85 minutes.\n"
    )
    return text_payload.encode("utf-8")


@dataclass(frozen=True, slots=True)
class DocumentCase:
    """인제스트할 실제 문서와 형식별 검증 계약을 정의한다."""

    name: str
    file_idx: int
    file_name: str
    file_type: str
    content_type: str
    parser_type: str
    parser_version: str
    answer_token: str
    search_query: str
    document_bytes: bytes
    expected_source_metadata: Mapping[str, object]

    @property
    def sha256(self) -> str:
        """현재 테스트 실행에서 사용하는 원본 문서 SHA-256을 반환한다."""

        return hashlib.sha256(self.document_bytes).hexdigest()

    @property
    def download_url(self) -> str:
        """Presigned GET URL 역할의 결정적인 HTTPS URL을 반환한다."""

        return (
            f"https://files.e2e.invalid/files/{self.file_idx}.{self.file_type}"
            f"?X-Amz-Signature=e2e-{self.file_idx}"
        )

    @property
    def manifest(self) -> dict[str, object]:
        """Backend manifest와 POST /ingest에 사용할 동일한 요청 본문을 만든다."""

        return {
            "file_idx": self.file_idx,
            "user_idx": _TEST_USER_IDX,
            "folder_idx": _TEST_FOLDER_IDX,
            "file_name": self.file_name,
            "file_type": self.file_type,
            "download_url": self.download_url,
            "url_expires_in": 900,
        }


_DOCX: Final[DocumentCase] = DocumentCase(
    name="docx",
    file_idx=_DOCX_FILE_IDX,
    file_name="jipsa-e2e-delta.docx",
    file_type="docx",
    content_type=("application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    # 운영 Factory는 이미지 유무와 관계없이 OCR-aware DOCX 파서를 선택한다.
    parser_type="DOCX_HYBRID_OCR",
    parser_version=_OCR_PARSER_VERSION,
    answer_token="DOCX-DELTA-52",
    search_query="DOCX 문서에 기록된 exact verification code는 무엇인가요?",
    document_bytes=_build_docx(),
    expected_source_metadata={
        "location_kind": "docx_block",
        "section_index": 1,
        "paragraph_index": 2,
        "section_title": "JIPSA DOCX E2E",
    },
)

_PPTX: Final[DocumentCase] = DocumentCase(
    name="pptx",
    file_idx=_PPTX_FILE_IDX,
    file_name="jipsa-e2e-echo.pptx",
    file_type="pptx",
    content_type=("application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    # 텍스트 도형만 포함한 Fixture도 운영상 Hybrid OCR 파서로 처리된다.
    parser_type="PPTX_HYBRID_OCR",
    parser_version=_OCR_PARSER_VERSION,
    answer_token="PPTX-ECHO-63",
    search_query="PPTX 문서에 기록된 exact verification code는 무엇인가요?",
    document_bytes=_build_pptx(),
    expected_source_metadata={
        "location_kind": "pptx_shape",
        "slide_number": 1,
    },
)

_XLSX: Final[DocumentCase] = DocumentCase(
    name="xlsx",
    file_idx=_XLSX_FILE_IDX,
    file_name="jipsa-e2e-foxtrot.xlsx",
    file_type="xlsx",
    content_type=("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    # 기본 셀 파싱은 유지되지만 최종 Parser 식별자는 이미지·차트 OCR까지 포괄한다.
    parser_type="XLSX_HYBRID_OCR",
    parser_version=_OCR_PARSER_VERSION,
    answer_token="XLSX-FOXTROT-74",
    search_query="XLSX 문서에 기록된 exact verification code는 무엇인가요?",
    document_bytes=_build_xlsx(),
    expected_source_metadata={
        "location_kind": "xlsx_cell_range",
        "sheet_name": "E2E",
        "row_number": 2,
        "cell_range": "A2:B2",
    },
)

_TXT: Final[DocumentCase] = DocumentCase(
    name="txt",
    file_idx=_TXT_FILE_IDX,
    file_name="jipsa-e2e-golf.txt",
    file_type="txt",
    content_type="text/plain; charset=utf-8",
    # TXT는 이미지 추출 대상이 아니므로 기존 텍스트 파서 계약을 그대로 유지한다.
    parser_type="TXT_TEXT",
    parser_version="1.1.0",
    answer_token="TXT-GOLF-85",
    search_query="TXT 문서에 기록된 exact verification code는 무엇인가요?",
    document_bytes=_build_txt(),
    expected_source_metadata={
        "location_kind": "txt_line",
        "line_number": 2,
    },
)

_DOCUMENTS: Final[tuple[DocumentCase, ...]] = (
    _DOCX,
    _PPTX,
    _XLSX,
    _TXT,
)
_DOCUMENTS_BY_IDX: Final[dict[int, DocumentCase]] = {case.file_idx: case for case in _DOCUMENTS}
_DOCUMENT_IDS: Final[tuple[str, ...]] = tuple(case.name for case in _DOCUMENTS)


# ============================================================
# 동적 JSON·DB·Qdrant 값 검증 도우미
# ============================================================


def _object(value: object, label: str) -> dict[str, object]:
    """동적 값을 문자열 key를 가진 JSON 객체로 좁힌다."""

    if not isinstance(value, dict):
        raise AssertionError(f"{label} must be a JSON object.")
    if any(not isinstance(key, str) for key in value):
        raise AssertionError(f"{label} must contain only string keys.")
    return cast(dict[str, object], value)


def _objects(
    mapping: Mapping[str, object],
    key: str,
) -> list[dict[str, object]]:
    """매핑에서 JSON 객체 배열을 읽는다."""

    value = mapping.get(key)
    if not isinstance(value, list):
        raise AssertionError(f"{key} must be a JSON array.")
    return [_object(item, key) for item in cast(list[object], value)]


def _str(mapping: Mapping[str, object], key: str) -> str:
    """매핑에서 비어 있지 않은 문자열을 읽는다."""

    value = mapping.get(key)
    if not isinstance(value, str) or not value:
        raise AssertionError(f"{key} must be a non-empty string.")
    return value


def _int(mapping: Mapping[str, object], key: str) -> int:
    """매핑에서 bool이 아닌 정수를 읽는다."""

    value = mapping.get(key)
    if isinstance(value, bool) or not isinstance(value, int):
        raise AssertionError(f"{key} must be an integer.")
    return value


def _optional_int(mapping: Mapping[str, object], key: str) -> int | None:
    """매핑에서 null 또는 bool이 아닌 정수를 읽는다."""

    value = mapping.get(key)
    if value is None:
        return None
    if isinstance(value, bool) or not isinstance(value, int):
        raise AssertionError(f"{key} must be an integer or null.")
    return value


def _bool(mapping: Mapping[str, object], key: str) -> bool:
    """매핑에서 JSON boolean 값을 읽는다."""

    value = mapping.get(key)
    if not isinstance(value, bool):
        raise AssertionError(f"{key} must be a boolean.")
    return value


def _db_bool(mapping: Mapping[str, object], key: str) -> bool:
    """MySQL/MariaDB의 bool 또는 0·1 값을 Python bool로 정규화한다."""

    value = mapping.get(key)
    if isinstance(value, bool):
        return value
    if isinstance(value, int) and value in {0, 1}:
        return bool(value)
    raise AssertionError(f"{key} must be a database boolean.")


def _source_metadata(mapping: Mapping[str, object]) -> dict[str, object]:
    """DB 또는 Qdrant에서 읽은 source_metadata를 JSON 객체로 정규화한다."""

    value = mapping.get("source_metadata")

    if isinstance(value, dict):
        return _object(value, "source_metadata")

    if isinstance(value, bytes):
        value = value.decode("utf-8")

    if isinstance(value, str):
        return _object(json.loads(value), "source_metadata")

    raise AssertionError("source_metadata must be a JSON object or JSON string.")


def _assert_metadata_contains(
    actual: Mapping[str, object],
    expected: Mapping[str, object],
) -> None:
    """형식별 핵심 위치 메타데이터가 값까지 동일하게 보존됐는지 검증한다."""

    for key, expected_value in expected.items():
        assert actual.get(key) == expected_value, (
            f"source_metadata[{key!r}] mismatch: "
            f"expected={expected_value!r}, actual={actual.get(key)!r}"
        )


# ============================================================
# Backend manifest·완료 콜백 및 다운로드 계약
# ============================================================


@dataclass(slots=True)
class BackendRecorder:
    """최신 manifest를 반환하고 실제 완료 콜백 payload를 기록한다."""

    settings: Settings
    cases: Mapping[int, DocumentCase]
    manifest_requests: list[int] = field(default_factory=list)
    callbacks: dict[int, list[dict[str, object]]] = field(default_factory=dict)

    async def handle(self, request: httpx2.Request) -> httpx2.Response:
        """Backend 내부 manifest와 ingest-complete 두 경로만 처리한다."""

        path_match = _BACKEND_PATH_PATTERN.fullmatch(request.url.path)
        if path_match is None:
            return httpx2.Response(status_code=404)

        file_idx = int(path_match.group("file_idx"))
        operation = path_match.group("operation")
        case = self.cases.get(file_idx)

        if case is None:
            return httpx2.Response(status_code=404)

        internal_token = self.settings.internal_token
        if internal_token is None:
            raise AssertionError("INTERNAL_TOKEN must be configured for E2E.")

        assert request.headers["X-Internal-Token"] == (internal_token.get_secret_value())

        if operation == "manifest":
            assert request.method == "GET"
            self.manifest_requests.append(file_idx)
            return httpx2.Response(status_code=200, json=case.manifest)

        assert operation == "ingest-complete"
        assert request.method == "POST"

        payload = _object(
            json.loads(request.content.decode("utf-8")),
            "ingest-complete payload",
        )
        self.callbacks.setdefault(file_idx, []).append(payload)
        return httpx2.Response(status_code=204)

    def callback(self, file_idx: int) -> dict[str, object]:
        """파일별 성공 콜백이 정확히 한 번 전송되었는지 확인한다."""

        payloads = self.callbacks.get(file_idx, [])
        assert len(payloads) == 1, (
            f"file_idx={file_idx} expected one callback, received {len(payloads)}"
        )
        payload = payloads[0]
        assert _bool(payload, "success") is True
        return payload


@dataclass(frozen=True, slots=True)
class DownloadContract:
    """HttpFileDownloader에 실제 형식별 문서 ByteStream을 제공한다."""

    cases: Mapping[int, DocumentCase]

    async def handle(self, request: httpx2.Request) -> httpx2.Response:
        """허용된 E2E 문서 GET 요청만 처리하고 원본 바이트를 그대로 반환한다."""

        path_match = _DOWNLOAD_PATH_PATTERN.fullmatch(request.url.path)
        if path_match is None:
            return httpx2.Response(status_code=404)

        file_idx = int(path_match.group("file_idx"))
        extension = path_match.group("extension")
        case = self.cases.get(file_idx)

        if case is None or extension != case.file_type:
            return httpx2.Response(status_code=404)

        assert request.method == "GET"
        assert request.headers["accept-encoding"] == "identity"

        return httpx2.Response(
            status_code=200,
            headers={
                "Content-Type": case.content_type,
                "Content-Length": str(len(case.document_bytes)),
            },
            stream=httpx2.ByteStream(case.document_bytes),
        )


# ============================================================
# Local RAG DB 조회 및 정리
# ============================================================


@dataclass(frozen=True, slots=True)
class DatabaseState:
    """한 파일의 활성 문서, 청크와 최신 색인 실행 상태."""

    document: Mapping[str, object]
    chunks: tuple[Mapping[str, object], ...]
    latest_run: Mapping[str, object]


@dataclass(frozen=True, slots=True)
class E2eRuntime:
    """네 형식의 실제 인제스트 완료 후 모듈 테스트가 공유할 상태."""

    client: TestClient
    settings: Settings
    recorder: BackendRecorder
    responses: Mapping[int, Mapping[str, object]]


def _db_engine(settings: Settings) -> AsyncEngine:
    """검증과 정리에만 사용하는 독립적인 비동기 DB 엔진을 만든다."""

    return create_async_engine(settings.database_url, pool_pre_ping=True)


async def _database_state(
    settings: Settings,
    file_idx: int,
) -> DatabaseState:
    """지정 파일의 실제 Local RAG 문서·청크·최신 실행 상태를 조회한다."""

    engine = _db_engine(settings)

    try:
        async with engine.connect() as connection:
            document_rows = (
                (
                    await connection.execute(
                        text(
                            """
                            SELECT
                                `RAG_Document_IDX` AS `rag_document_idx`,
                                `File_IDX` AS `file_idx`,
                                `Users_IDX` AS `users_idx`,
                                `Folder_IDX` AS `folder_idx`,
                                `File_Name` AS `file_name`,
                                `File_Type` AS `file_type`,
                                `File_Hash` AS `file_hash`,
                                `Index_Version` AS `index_version`,
                                `Parse_Status` AS `parse_status`,
                                `Index_Status` AS `index_status`,
                                `Chunk_Count` AS `chunk_count`,
                                `Parser_Type` AS `parser_type`,
                                `Parser_Version` AS `parser_version`,
                                `Embedding_Model` AS `embedding_model`,
                                (`Deleted_At` IS NOT NULL) AS `is_deleted`
                            FROM `RAG_Document`
                            WHERE `Users_IDX` = :users_idx
                              AND `File_IDX` = :file_idx
                              AND `Deleted_At` IS NULL
                            ORDER BY `RAG_Document_IDX`
                            """
                        ),
                        {
                            "users_idx": _TEST_USER_IDX,
                            "file_idx": file_idx,
                        },
                    )
                )
                .mappings()
                .all()
            )
            assert len(document_rows) == 1
            document = cast(Mapping[str, object], document_rows[0])
            rag_document_idx = _int(document, "rag_document_idx")

            chunk_rows = (
                (
                    await connection.execute(
                        text(
                            """
                            SELECT
                                `Chunk_ID` AS `chunk_id`,
                                `RAG_Document_IDX` AS `rag_document_idx`,
                                `File_IDX` AS `file_idx`,
                                `Users_IDX` AS `users_idx`,
                                `Folder_IDX` AS `folder_idx`,
                                `Chunk_Index` AS `chunk_index`,
                                `Content` AS `content`,
                                `Token_Count` AS `token_count`,
                                `Page` AS `page`,
                                `Slide_No` AS `slide_no`,
                                `Sheet_Name` AS `sheet_name`,
                                `Section_Title` AS `section_title`,
                                `Start_Offset` AS `start_offset`,
                                `End_Offset` AS `end_offset`,
                                `Content_Hash` AS `content_hash`,
                                `Embedding_Model` AS `embedding_model`,
                                `Index_Version` AS `index_version`,
                                `Source_Metadata` AS `source_metadata`
                            FROM `RAG_Chunk`
                            WHERE `RAG_Document_IDX` = :rag_document_idx
                            ORDER BY `Chunk_Index`
                            """
                        ),
                        {"rag_document_idx": rag_document_idx},
                    )
                )
                .mappings()
                .all()
            )

            run_rows = (
                (
                    await connection.execute(
                        text(
                            """
                            SELECT
                                `RAG_Index_Run_IDX` AS `rag_index_run_idx`,
                                `RAG_Document_IDX` AS `rag_document_idx`,
                                `File_IDX` AS `file_idx`,
                                `Users_IDX` AS `users_idx`,
                                `Run_Type` AS `run_type`,
                                `Status` AS `status`,
                                `Parser_Type` AS `parser_type`,
                                `Parser_Version` AS `parser_version`,
                                `Embedding_Model` AS `embedding_model`,
                                `Chunk_Count` AS `chunk_count`,
                                (`Finished_At` IS NOT NULL) AS `is_finished`,
                                `Error_Message` AS `error_message`
                            FROM `RAG_Index_Run`
                            WHERE `RAG_Document_IDX` = :rag_document_idx
                            ORDER BY `RAG_Index_Run_IDX` DESC
                            LIMIT 1
                            """
                        ),
                        {"rag_document_idx": rag_document_idx},
                    )
                )
                .mappings()
                .all()
            )
            assert len(run_rows) == 1

            return DatabaseState(
                document=document,
                chunks=tuple(cast(Mapping[str, object], row) for row in chunk_rows),
                latest_run=cast(Mapping[str, object], run_rows[0]),
            )
    finally:
        await engine.dispose()


async def _cleanup_database(settings: Settings) -> None:
    """E2E 전용 파일 범위의 실행·청크·문서를 FK 역순으로 삭제한다."""

    engine = _db_engine(settings)
    parameters = {
        "users_idx": _TEST_USER_IDX,
        "file_idx_min": _FILE_IDX_MIN,
        "file_idx_max": _FILE_IDX_MAX,
    }

    try:
        async with engine.begin() as connection:
            # 자식 테이블이 RAG_Document를 참조하므로 실행 이력과 청크를
            # 먼저 삭제한 뒤 문서 행을 삭제한다.
            for statement in (
                text(
                    """
                    DELETE FROM `RAG_Index_Run`
                    WHERE `Users_IDX` = :users_idx
                      AND `File_IDX`
                          BETWEEN :file_idx_min AND :file_idx_max
                    """
                ),
                text(
                    """
                    DELETE FROM `RAG_Chunk`
                    WHERE `Users_IDX` = :users_idx
                      AND `File_IDX`
                          BETWEEN :file_idx_min AND :file_idx_max
                    """
                ),
                text(
                    """
                    DELETE FROM `RAG_Document`
                    WHERE `Users_IDX` = :users_idx
                      AND `File_IDX`
                          BETWEEN :file_idx_min AND :file_idx_max
                    """
                ),
            ):
                await connection.execute(statement, parameters)
    finally:
        await engine.dispose()


# ============================================================
# Qdrant 조회 및 정리
# ============================================================


def _qdrant_client(settings: Settings) -> AsyncQdrantClient:
    """현재 테스트 환경의 실제 Qdrant 클라이언트를 생성한다."""

    api_key = (
        settings.qdrant_api_key.get_secret_value() if settings.qdrant_api_key is not None else None
    )
    return AsyncQdrantClient(
        url=settings.qdrant_url,
        grpc_port=settings.qdrant_grpc_port,
        prefer_grpc=settings.qdrant_prefer_grpc,
        api_key=api_key,
        timeout=max(1, ceil(settings.qdrant_timeout_seconds)),
    )


def _scope_filter(
    file_idxs: Sequence[int],
    *,
    active_only: bool,
) -> models.Filter:
    """E2E 사용자·파일 범위를 고정하고 선택적으로 활성 조건을 추가한다."""

    if active_only:
        return models.Filter(
            must=[
                models.FieldCondition(
                    key="users_idx",
                    match=models.MatchValue(value=_TEST_USER_IDX),
                ),
                models.FieldCondition(
                    key="file_idx",
                    match=models.MatchAny(any=list(file_idxs)),
                ),
                models.FieldCondition(
                    key="is_active",
                    match=models.MatchValue(value=True),
                ),
            ]
        )

    return models.Filter(
        must=[
            models.FieldCondition(
                key="users_idx",
                match=models.MatchValue(value=_TEST_USER_IDX),
            ),
            models.FieldCondition(
                key="file_idx",
                match=models.MatchAny(any=list(file_idxs)),
            ),
        ]
    )


async def _active_points(
    settings: Settings,
    file_idx: int,
) -> tuple[models.Record, ...]:
    """지정 파일의 활성 Point와 payload·vector를 실제 Qdrant에서 읽는다."""

    client = _qdrant_client(settings)
    try:
        points, next_offset = await client.scroll(
            collection_name=settings.qdrant_collection,
            scroll_filter=_scope_filter((file_idx,), active_only=True),
            limit=256,
            with_payload=True,
            with_vectors=True,
        )
        assert next_offset is None
        return tuple(points)
    finally:
        await client.close()


async def _cleanup_qdrant(settings: Settings) -> None:
    """E2E 범위의 활성·비활성 Point를 payload filter로 삭제한다."""

    client = _qdrant_client(settings)
    try:
        if not await client.collection_exists(settings.qdrant_collection):
            return
        await client.delete(
            collection_name=settings.qdrant_collection,
            points_selector=models.FilterSelector(
                filter=_scope_filter(_FILE_IDXS, active_only=False)
            ),
            wait=True,
        )
    finally:
        await client.close()


async def _cleanup(settings: Settings) -> None:
    """Qdrant 복제 데이터와 Local RAG 원본 데이터를 순서대로 정리한다."""

    await _cleanup_qdrant(settings)
    await _cleanup_database(settings)


# ============================================================
# 실제 E2E 인프라 설정 사전 검증
# ============================================================


def _uses_test_only_hostname(url: str) -> bool:
    """URL이 ``*.test`` 전용 가짜 호스트를 사용하는지 확인한다.

    일반 단위·통합 테스트의 ``.env.test``에는 실제 네트워크 호출을 막기 위해
    ``qdrant.test``와 ``embedding.test`` 같은 예약 도메인을 사용한다.
    이 주소는 실제 E2E에서 접속해야 하는 Local Qdrant·CUDA TEI 주소가 아니며,
    Windows DNS에서 해석되지 않는 것이 정상이다.
    """

    hostname = urlsplit(url).hostname
    if hostname is None:
        return True

    normalized_hostname = hostname.rstrip(".").lower()
    return normalized_hostname == "test" or normalized_hostname.endswith(".test")


def _validate_real_e2e_infrastructure_settings(settings: Settings) -> None:
    """실제 E2E가 Mock 전용 ``.env.test`` 주소로 실행되는 것을 차단한다."""

    invalid_settings: list[str] = []

    if _uses_test_only_hostname(str(settings.qdrant_url)):
        invalid_settings.append("JIPSA_RAG_QDRANT_URL")

    if _uses_test_only_hostname(str(settings.embedding_base_url)):
        invalid_settings.append("JIPSA_RAG_EMBEDDING_BASE_URL")

    if not invalid_settings:
        return

    invalid_setting_names = ", ".join(invalid_settings)
    pytest.fail(
        "Real non-PDF multiformat RAG E2E received .env.test mock-only "
        f"infrastructure settings: {invalid_setting_names}. "
        "Load .env.local into the current process, keep "
        "JIPSA_RAG_APP_ENV=test, start Local Qdrant and CUDA TEI, and then "
        "run this test. Do not replace the intentional qdrant.test or "
        "embedding.test values in .env.test.",
        pytrace=False,
    )


# ============================================================
# 실제 네 형식 인제스트 공통 Fixture
# ============================================================


@pytest.fixture(scope="module", autouse=True)
def require_e2e_opt_in() -> None:
    """일반 테스트 실행에서 실제 Claude 및 Local 인프라 호출을 방지한다."""

    if os.getenv(_RUN_ENV) != "1":
        pytest.skip(f"Set {_RUN_ENV}=1 to run real non-PDF multiformat RAG E2E tests.")


@pytest.fixture(scope="module")
def e2e_runtime(
    tmp_path_factory: pytest.TempPathFactory,
) -> Iterator[E2eRuntime]:
    """DOCX·PPTX·XLSX·TXT를 실제로 인제스트하고 모듈 동안 유지한다."""

    settings = get_settings()

    if settings.app_env != "test":
        pytest.fail(
            "Real E2E cleanup is allowed only when JIPSA_RAG_APP_ENV=test.",
            pytrace=False,
        )

    # 실제 네트워크 호출이나 E2E 전용 데이터 삭제보다 먼저 현재 Settings가
    # .env.local의 실제 인프라 주소를 사용하고 있는지 확인한다.
    _validate_real_e2e_infrastructure_settings(settings)

    get_generation_settings.cache_clear()
    try:
        get_generation_settings()
    except ValidationError as error:
        pytest.fail(
            f"A valid Anthropic API key/model is required for real E2E: {type(error).__name__}",
            pytrace=False,
        )

    ingest_token = settings.rag_ingest_token
    if ingest_token is None:
        pytest.fail(
            "RAG_INGEST_TOKEN is required for real E2E.",
            pytrace=False,
        )
    if settings.internal_token is None:
        pytest.fail(
            "INTERNAL_TOKEN is required for real E2E.",
            pytrace=False,
        )

    # Backend/S3 역할의 HTTP 경계만 테스트 대역으로 교체한다. DB, TEI,
    # Qdrant 및 Claude 관련 설정은 실제 E2E 프로세스 값을 그대로 사용한다.
    #
    # file_download_allowed_host_suffixes는 Settings 내부에서 쉼표로 구분된
    # 원시 문자열로 보관된다. model_copy(update=...)는 필드 타입을 재검증하지
    # 않으므로 반드시 Settings 계약과 동일한 문자열을 넣는다.
    http_settings = settings.model_copy(
        update={
            "app_server_base_url": "https://backend.e2e.invalid",
            "app_server_max_attempts": 1,
            "app_server_retry_initial_delay_seconds": 0.0,
            "app_server_retry_max_delay_seconds": 0.0,
            "file_download_allowed_host_suffixes": ".e2e.invalid",
        }
    )

    recorder = BackendRecorder(
        settings=http_settings,
        cases=_DOCUMENTS_BY_IDX,
    )
    download_contract = DownloadContract(cases=_DOCUMENTS_BY_IDX)
    backend_client = ApplicationServerIngestClient(
        http_settings,
        transport=httpx2.MockTransport(recorder.handle),
    )
    downloader = HttpFileDownloader(
        http_settings,
        transport=httpx2.MockTransport(download_contract.handle),
        temp_directory=Path(tmp_path_factory.mktemp("real-multiformat-rag-e2e")),
    )

    def backend_dependency() -> ApplicationServerIngestClient:
        """POST /ingest에 Backend HTTP 계약 대역만 주입한다."""

        return backend_client

    def downloader_dependency() -> HttpFileDownloader:
        """실제 downloader에 형식별 응답 transport만 주입한다."""

        return downloader

    app.dependency_overrides[get_application_server_ingest_client] = backend_dependency
    app.dependency_overrides[get_file_downloader] = downloader_dependency

    # 최초 정리가 성공하기 전에 Fixture가 실패하면 새 E2E 데이터는 생성되지
    # 않는다. 종료 정리는 실제 인제스트 가능 상태에 진입했을 때만 수행한다.
    initial_cleanup_completed = False

    try:
        asyncio.run(_cleanup(settings))
        initial_cleanup_completed = True

        with TestClient(app) as client:
            client.headers["X-Internal-Token"] = ingest_token.get_secret_value()
            responses: dict[int, Mapping[str, object]] = {}

            for case in _DOCUMENTS:
                response = client.post("/ingest", json=case.manifest)
                assert response.status_code == 200, (
                    f"{case.name} ingest failed: "
                    f"status={response.status_code}, body={response.text}"
                )

                body = _object(response.json(), "POST /ingest response")
                assert _bool(body, "success") is True
                assert _str(body, "code") == "FILE_INDEXING_COMPLETED"
                responses[case.file_idx] = _object(
                    body.get("data"),
                    "POST /ingest response data",
                )

            yield E2eRuntime(
                client=client,
                settings=settings,
                recorder=recorder,
                responses=responses,
            )
    finally:
        try:
            if initial_cleanup_completed:
                asyncio.run(_cleanup(settings))
        finally:
            app.dependency_overrides.pop(
                get_application_server_ingest_client,
                None,
            )
            app.dependency_overrides.pop(get_file_downloader, None)
            get_generation_settings.cache_clear()


# ============================================================
# 1. 실제 문서 Fixture 자체 계약
# ============================================================


def _extract_fixture_text(case: DocumentCase) -> str:
    """운영 파서와 독립적인 라이브러리 경로로 Fixture 원문을 읽는다."""

    if case.file_type == "docx":
        document = Document(BytesIO(case.document_bytes))
        paragraphs = [paragraph.text for paragraph in document.paragraphs]
        table_values = [
            cell.text for table in document.tables for row in table.rows for cell in row.cells
        ]
        return "\n".join((*paragraphs, *table_values))

    if case.file_type == "pptx":
        presentation = Presentation(BytesIO(case.document_bytes))
        values: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                shape_text = getattr(shape, "text", None)
                if isinstance(shape_text, str) and shape_text:
                    values.append(shape_text)

                shape_table = getattr(shape, "table", None)
                if shape_table is not None:
                    values.extend(cell.text for row in shape_table.rows for cell in row.cells)
        return "\n".join(values)

    if case.file_type == "xlsx":
        workbook = load_workbook(
            BytesIO(case.document_bytes),
            read_only=True,
            data_only=False,
        )
        try:
            values = [
                str(cell.value)
                for worksheet in workbook.worksheets
                for row in worksheet.iter_rows()
                for cell in row
                if cell.value is not None
            ]
            return "\n".join(values)
        finally:
            workbook.close()

    if case.file_type == "txt":
        return case.document_bytes.decode("utf-8")

    raise AssertionError(f"Unexpected E2E document type: {case.file_type}")


@pytest.mark.parametrize("case", _DOCUMENTS, ids=_DOCUMENT_IDS)
def test_real_multiformat_fixture_contract(case: DocumentCase) -> None:
    """실제 문서 바이트와 고유 검색 토큰이 형식별로 유효한지 확인한다."""

    assert case.document_bytes
    assert _SHA256_PATTERN.fullmatch(case.sha256)
    assert case.answer_token in _extract_fixture_text(case)

    if case.file_type in {"docx", "pptx", "xlsx"}:
        assert case.document_bytes.startswith(b"PK")
    else:
        assert case.file_type == "txt"
        assert not case.document_bytes.startswith(b"PK")


# ============================================================
# 2. 실제 인제스트·청킹·완료 콜백 계약
# ============================================================


@pytest.mark.parametrize("case", _DOCUMENTS, ids=_DOCUMENT_IDS)
def test_actual_multiformat_ingest_chunking_and_callback(
    e2e_runtime: E2eRuntime,
    case: DocumentCase,
) -> None:
    """형식별 실제 파싱·청킹 결과가 API와 Backend 콜백에 보존되는지 검증한다."""

    response_data = e2e_runtime.responses[case.file_idx]
    callback = e2e_runtime.recorder.callback(case.file_idx)

    assert e2e_runtime.recorder.manifest_requests.count(case.file_idx) == 1
    assert _int(response_data, "file_idx") == case.file_idx
    assert _int(response_data, "user_idx") == _TEST_USER_IDX
    assert _optional_int(response_data, "folder_idx") == _TEST_FOLDER_IDX
    assert _str(response_data, "file_name") == case.file_name
    assert _str(response_data, "file_type") == case.file_type
    assert _int(response_data, "file_size_bytes") == len(case.document_bytes)
    assert _int(response_data, "page_count") > 0
    assert _int(response_data, "text_unit_count") > 0
    assert _int(response_data, "chunk_count") > 0
    assert _str(response_data, "embedding_model") == (e2e_runtime.settings.embedding_model)
    assert _int(response_data, "embedding_dim") == (e2e_runtime.settings.embedding_dim)
    assert _str(response_data, "processing_status") == "INDEXED"

    callback_chunks = _objects(callback, "chunks")
    assert _int(callback, "index_version") == _INDEX_VERSION
    assert _int(callback, "chunk_count") == len(callback_chunks)
    assert len(callback_chunks) == _int(response_data, "chunk_count")

    token_chunks = [
        chunk for chunk in callback_chunks if case.answer_token in _str(chunk, "content")
    ]
    assert token_chunks, f"{case.answer_token} was not preserved in callback chunks."

    for chunk in callback_chunks:
        assert _str(chunk, "chunk_id")
        assert _int(chunk, "chunk_index") >= 0
        assert _SHA256_PATTERN.fullmatch(_str(chunk, "content_hash"))

    for token_chunk in token_chunks:
        metadata = _object(
            token_chunk.get("source_metadata"),
            "callback source_metadata",
        )
        _assert_metadata_contains(metadata, case.expected_source_metadata)
        assert metadata.get("chunking_strategy") == "STRUCTURED_DOCUMENT"
        assert metadata.get("chunking_strategy_version") == "1.0.0"


# ============================================================
# 3. Local RAG DB 문서·청크·색인 실행 상태
# ============================================================


@pytest.mark.parametrize("case", _DOCUMENTS, ids=_DOCUMENT_IDS)
def test_local_rag_multiformat_document_chunk_and_index_state(
    e2e_runtime: E2eRuntime,
    case: DocumentCase,
) -> None:
    """Local RAG DB에 형식·파서·전체 위치 메타데이터가 저장되는지 검증한다."""

    state = asyncio.run(_database_state(e2e_runtime.settings, case.file_idx))
    document = state.document
    latest_run = state.latest_run

    assert _int(document, "file_idx") == case.file_idx
    assert _int(document, "users_idx") == _TEST_USER_IDX
    assert _optional_int(document, "folder_idx") == _TEST_FOLDER_IDX
    assert _str(document, "file_name") == case.file_name
    assert _str(document, "file_type") == case.file_type.upper()
    assert _str(document, "file_hash") == case.sha256
    assert _int(document, "index_version") == _INDEX_VERSION
    assert _str(document, "parse_status") == "PARSED"
    assert _str(document, "index_status") == "INDEXED"
    assert _str(document, "parser_type") == case.parser_type
    assert _str(document, "parser_version") == case.parser_version
    assert _str(document, "embedding_model") == (e2e_runtime.settings.embedding_model)
    assert _db_bool(document, "is_deleted") is False
    assert _int(document, "chunk_count") == len(state.chunks)
    assert state.chunks

    assert tuple(_int(chunk, "chunk_index") for chunk in state.chunks) == tuple(
        range(len(state.chunks))
    )
    assert case.answer_token in "\n".join(_str(chunk, "content") for chunk in state.chunks)

    token_chunks = [chunk for chunk in state.chunks if case.answer_token in _str(chunk, "content")]
    assert token_chunks

    for chunk in state.chunks:
        assert _int(chunk, "file_idx") == case.file_idx
        assert _int(chunk, "users_idx") == _TEST_USER_IDX
        assert _optional_int(chunk, "folder_idx") == _TEST_FOLDER_IDX
        assert _int(chunk, "start_offset") >= 0
        assert _int(chunk, "end_offset") > _int(chunk, "start_offset")
        assert _SHA256_PATTERN.fullmatch(_str(chunk, "content_hash"))
        assert _str(chunk, "embedding_model") == (e2e_runtime.settings.embedding_model)
        assert _int(chunk, "index_version") == _INDEX_VERSION

    for token_chunk in token_chunks:
        metadata = _source_metadata(token_chunk)
        _assert_metadata_contains(metadata, case.expected_source_metadata)

        if case.file_type == "docx":
            assert _str(token_chunk, "section_title") == "JIPSA DOCX E2E"
        elif case.file_type == "pptx":
            assert _optional_int(token_chunk, "slide_no") == 1
        elif case.file_type == "xlsx":
            assert _str(token_chunk, "sheet_name") == "E2E"
        elif case.file_type == "txt":
            assert _optional_int(token_chunk, "page") is None
            assert _optional_int(token_chunk, "slide_no") is None
            assert token_chunk.get("sheet_name") is None

    assert _int(latest_run, "rag_document_idx") == _int(
        document,
        "rag_document_idx",
    )
    assert _int(latest_run, "file_idx") == case.file_idx
    assert _int(latest_run, "users_idx") == _TEST_USER_IDX
    assert _str(latest_run, "run_type") == "FULL"
    assert _str(latest_run, "status") == "SUCCESS"
    assert _str(latest_run, "parser_type") == case.parser_type
    assert _str(latest_run, "parser_version") == case.parser_version
    assert _str(latest_run, "embedding_model") == (e2e_runtime.settings.embedding_model)
    assert _int(latest_run, "chunk_count") == len(state.chunks)
    assert _db_bool(latest_run, "is_finished") is True
    assert latest_run.get("error_message") is None


# ============================================================
# 4. Qdrant 벡터 및 형식별 payload 계약
# ============================================================


@pytest.mark.parametrize("case", _DOCUMENTS, ids=_DOCUMENT_IDS)
def test_qdrant_multiformat_active_points_and_payload(
    e2e_runtime: E2eRuntime,
    case: DocumentCase,
) -> None:
    """Local RAG 청크와 Qdrant Point·vector·source_metadata를 대조한다."""

    state = asyncio.run(_database_state(e2e_runtime.settings, case.file_idx))
    points = asyncio.run(_active_points(e2e_runtime.settings, case.file_idx))
    chunks_by_id = {_str(chunk, "chunk_id"): chunk for chunk in state.chunks}

    assert len(points) == len(chunks_by_id)
    assert {str(point.id) for point in points} == set(chunks_by_id)

    token_payloads: list[Mapping[str, object]] = []

    for point in points:
        assert point.payload is not None
        payload = cast(Mapping[str, object], point.payload)
        point_id = str(point.id)
        local_chunk = chunks_by_id[point_id]

        assert isinstance(point.vector, list)
        assert len(point.vector) == e2e_runtime.settings.embedding_dim
        assert _str(payload, "chunk_id") == point_id
        assert _int(payload, "file_idx") == case.file_idx
        assert _int(payload, "users_idx") == _TEST_USER_IDX
        assert _optional_int(payload, "folder_idx") == _TEST_FOLDER_IDX
        assert _str(payload, "file_name") == case.file_name
        assert _str(payload, "file_type") == case.file_type.upper()
        assert _str(payload, "file_hash") == case.sha256
        assert _str(payload, "parser_type") == case.parser_type
        assert _str(payload, "parser_version") == case.parser_version
        assert _str(payload, "content") == _str(local_chunk, "content")
        assert _str(payload, "content_hash") == _str(
            local_chunk,
            "content_hash",
        )
        assert _str(payload, "embedding_model") == (e2e_runtime.settings.embedding_model)
        assert _int(payload, "embedding_dim") == (e2e_runtime.settings.embedding_dim)
        assert _int(payload, "index_version") == _INDEX_VERSION
        assert _bool(payload, "is_active") is True
        assert _str(payload, "created_at")

        if case.answer_token in _str(payload, "content"):
            token_payloads.append(payload)

    assert token_payloads
    for payload in token_payloads:
        _assert_metadata_contains(
            _source_metadata(payload),
            case.expected_source_metadata,
        )


# ============================================================
# 5. 실제 TEI·Qdrant 형식별 검색
# ============================================================


@pytest.mark.parametrize("case", _DOCUMENTS, ids=_DOCUMENT_IDS)
def test_real_multiformat_chunk_search(
    e2e_runtime: E2eRuntime,
    case: DocumentCase,
) -> None:
    """각 형식이 실제 질의 임베딩과 Qdrant 검색 결과로 반환되는지 검증한다."""

    response = e2e_runtime.client.post(
        "/api/v1/chunks/search",
        json={
            "user_idx": _TEST_USER_IDX,
            "reference_file_idxs": [case.file_idx],
            "query": case.search_query,
            "top_k": 10,
            "score_threshold": None,
        },
    )
    assert response.status_code == 200, (
        f"{case.name} search failed: status={response.status_code}, body={response.text}"
    )

    body = _object(response.json(), "chunk search response")
    assert _bool(body, "success") is True
    assert _str(body, "code") == "CHUNK_SEARCH_COMPLETED"
    data = _object(body.get("data"), "chunk search response data")
    results = _objects(data, "results")

    assert results
    assert {_int(result, "file_idx") for result in results} == {case.file_idx}
    assert {_str(result, "file_type") for result in results} == {case.file_type}
    assert case.answer_token in "\n".join(_str(result, "content") for result in results)

    token_results = [result for result in results if case.answer_token in _str(result, "content")]
    assert token_results

    for result in token_results:
        if case.file_type == "docx":
            assert _str(result, "section_title") == "JIPSA DOCX E2E"
        elif case.file_type == "pptx":
            assert _optional_int(result, "slide_no") == 1
        elif case.file_type == "xlsx":
            assert _str(result, "sheet_name") == "E2E"
        elif case.file_type == "txt":
            assert _optional_int(result, "page") is None
            assert _optional_int(result, "slide_no") is None
            assert result.get("sheet_name") is None


# ============================================================
# 6. 네 형식 혼합 검색 기반 실제 Claude 답변
# ============================================================


def test_real_multiformat_mixed_document_claude_answer(
    e2e_runtime: E2eRuntime,
) -> None:
    """DOCX·PPTX·XLSX·TXT 근거를 한 답변에서 모두 검색·인용하는지 검증한다."""

    response = e2e_runtime.client.post(
        "/api/v1/rag/answers",
        json={
            "user_idx": _TEST_USER_IDX,
            "reference_file_idxs": list(_FILE_IDXS),
            "query": (
                "선택한 DOCX, PPTX, XLSX, TXT 문서 각각에 기록된 "
                "exact verification code를 파일 형식별로 원문 그대로 답하고, "
                "각 코드 문장 뒤에 해당 문서 출처를 반드시 인용해 주세요."
            ),
            # 네 문서에서 제목, 본문, 표가 여러 청크로 생성되더라도 모든
            # 관련 후보를 Claude 문맥에 포함할 수 있도록 API 최대값을 사용한다.
            "top_k": 20,
            "score_threshold": None,
        },
    )
    assert response.status_code == 200, (
        "mixed multiformat Claude answer failed: "
        f"status={response.status_code}, body={response.text}"
    )

    body = _object(response.json(), "RAG answer response")
    assert _bool(body, "success") is True
    assert _str(body, "code") == "RAG_ANSWER_COMPLETED"
    data = _object(body.get("data"), "RAG answer response data")
    answer = _str(data, "answer")

    assert _str(data, "status") == "answered"
    assert _str(data, "model").startswith("claude-")

    usage = _object(data.get("usage"), "RAG answer usage")
    assert _int(usage, "input_tokens") > 0
    assert _int(usage, "output_tokens") > 0

    for case in _DOCUMENTS:
        assert case.answer_token in answer

    sources = _objects(data, "sources")
    assert frozenset(_int(source, "file_idx") for source in sources) == frozenset(_FILE_IDXS)
    assert frozenset(_str(source, "file_type") for source in sources) == frozenset(
        case.file_type for case in _DOCUMENTS
    )

    cited_source_ids = tuple(dict.fromkeys(_SOURCE_PATTERN.findall(answer)))
    response_source_ids = tuple(_str(source, "source_id") for source in sources)
    assert cited_source_ids
    assert cited_source_ids == response_source_ids

    for source in sources:
        source_file_idx = _int(source, "file_idx")
        case = _DOCUMENTS_BY_IDX[source_file_idx]

        assert _str(source, "file_name") == case.file_name
        assert _str(source, "file_type") == case.file_type
        assert _str(source, "excerpt")

        if case.file_type == "docx":
            assert _str(source, "section_title") == "JIPSA DOCX E2E"
        elif case.file_type == "pptx":
            assert _optional_int(source, "slide_no") == 1
        elif case.file_type == "xlsx":
            assert _str(source, "sheet_name") == "E2E"
        elif case.file_type == "txt":
            assert _optional_int(source, "page") is None
            assert _optional_int(source, "slide_no") is None
            assert source.get("sheet_name") is None

        # 최종 API 출처의 Chunk_ID가 실제 Local RAG DB에 저장된 해당 파일
        # 청크 중 하나인지 확인하여 Claude 응답과 저장 원본의 연결을 검증한다.
        state = asyncio.run(_database_state(e2e_runtime.settings, source_file_idx))
        assert _str(source, "chunk_id") in {_str(chunk, "chunk_id") for chunk in state.chunks}
