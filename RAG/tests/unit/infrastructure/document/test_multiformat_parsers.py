"""PDF 외 지원 형식 파서의 핵심 추출 계약을 검증한다."""

from pathlib import Path

import pytest
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from openpyxl import Workbook  # type: ignore[import-untyped]
from openpyxl.worksheet.table import (  # type: ignore[import-untyped]
    Table,
    TableStyleInfo,
)
from pptx import Presentation
from pptx.util import Inches

from jipsa_rag.infrastructure.document.exceptions import InvalidDocumentError
from jipsa_rag.infrastructure.document.models import DocumentType
from jipsa_rag.infrastructure.document.parsers.docx import DocxDocumentParser
from jipsa_rag.infrastructure.document.parsers.pptx import PptxDocumentParser
from jipsa_rag.infrastructure.document.parsers.txt import TxtDocumentParser
from jipsa_rag.infrastructure.document.parsers.xlsx import XlsxDocumentParser


@pytest.mark.asyncio
async def test_docx_extracts_heading_paragraph_list_and_table(tmp_path: Path) -> None:
    path = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("프로젝트 개요", level=1)
    document.add_paragraph("일반 문단")
    document.add_paragraph("첫 번째 항목", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "이름"
    table.cell(0, 1).text = "값"
    table.cell(1, 0).text = "GPU"
    table.cell(1, 1).text = "CUDA 12.9"
    document.save(path)

    parsed = await DocxDocumentParser().parse(path)

    assert parsed.file_type is DocumentType.DOCX
    assert any(unit.source_metadata["unit_type"] == "heading" for unit in parsed.units)
    assert any(unit.source_metadata["unit_type"] == "list_item" for unit in parsed.units)
    assert any(unit.source_metadata["unit_type"] == "table" for unit in parsed.units)
    assert "CUDA 12.9" in parsed.text
    assert all("section_index" in unit.source_metadata for unit in parsed.units)


@pytest.mark.asyncio
async def test_docx_increments_section_metadata_after_section_break(tmp_path: Path) -> None:
    path = tmp_path / "sections.docx"
    document = Document()
    document.add_paragraph("첫 번째 섹션")
    document.add_section(WD_SECTION.NEW_PAGE)
    document.add_paragraph("두 번째 섹션")
    document.save(path)

    parsed = await DocxDocumentParser().parse(path)
    text_units = [unit for unit in parsed.units if unit.text]

    assert [unit.text for unit in text_units] == ["첫 번째 섹션", "두 번째 섹션"]
    assert [unit.source_metadata["section_index"] for unit in text_units] == [1, 2]
    assert parsed.document_metadata["section_count"] == 2


@pytest.mark.asyncio
async def test_docx_resolves_heading_and_list_from_base_styles(tmp_path: Path) -> None:
    path = tmp_path / "inherited-styles.docx"
    document = Document()

    custom_heading = document.styles.add_style(
        "Project Heading",
        WD_STYLE_TYPE.PARAGRAPH,
    )
    custom_heading.base_style = document.styles["Heading 2"]

    custom_list = document.styles.add_style(
        "Project List",
        WD_STYLE_TYPE.PARAGRAPH,
    )
    custom_list.base_style = document.styles["List Bullet"]

    document.add_paragraph("상속 제목", style=custom_heading)
    document.add_paragraph("상속 목록", style=custom_list)
    document.save(path)

    parsed = await DocxDocumentParser().parse(path)
    units_by_text = {unit.text: unit for unit in parsed.units if unit.text}

    assert units_by_text["상속 제목"].source_metadata["unit_type"] == "heading"
    assert units_by_text["상속 제목"].source_metadata["heading_level"] == 2
    assert units_by_text["상속 목록"].source_metadata["unit_type"] == "list_item"


@pytest.mark.asyncio
async def test_pptx_extracts_title_shape_table_and_notes(tmp_path: Path) -> None:
    path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title_shape = slide.shapes.title
    assert title_shape is not None
    title_shape.text = "인제스트 설계"
    slide.placeholders[1].text = "도형 본문"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(2))
    table_shape.table.cell(0, 0).text = "형식"
    table_shape.table.cell(0, 1).text = "지원"
    table_shape.table.cell(1, 0).text = "PPTX"
    table_shape.table.cell(1, 1).text = "예"
    slide.notes_slide.notes_text_frame.text = "발표자 노트"
    presentation.save(path)

    parsed = await PptxDocumentParser().parse(path)

    assert parsed.file_type is DocumentType.PPTX
    assert parsed.document_metadata["slide_count"] == 1
    assert {unit.source_metadata["unit_type"] for unit in parsed.units} >= {
        "title",
        "shape_text",
        "table",
        "speaker_notes",
    }
    assert all(unit.source_metadata["slide_number"] == 1 for unit in parsed.units)


@pytest.mark.asyncio
async def test_xlsx_extracts_rows_merged_ranges_tables_and_formula_metadata(tmp_path: Path) -> None:
    path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "요약"
    sheet.append(["항목", "값"])
    sheet.append(["문서 수", 5])
    sheet.append(["합계", "=SUM(B2:B2)"])
    sheet.merge_cells("A4:B4")
    sheet["A4"] = "병합 셀"
    table = Table(displayName="SummaryTable", ref="A1:B3")
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True)
    sheet.add_table(table)
    workbook.save(path)

    parsed = await XlsxDocumentParser().parse(path)

    assert parsed.file_type is DocumentType.XLSX
    assert parsed.document_metadata["sheet_names"] == ("요약",)
    assert "SummaryTable" in parsed.document_metadata["table_names"]
    formula_unit = next(unit for unit in parsed.units if unit.source_metadata["formula_cells"])
    assert formula_unit.source_metadata["formula_cells"] == ("B3",)
    assert any("A4:B4" in unit.source_metadata["merged_ranges"] for unit in parsed.units)


@pytest.mark.asyncio
async def test_ooxml_parser_rejects_another_ooxml_package_type(tmp_path: Path) -> None:
    path = tmp_path / "renamed.pptx"
    document = Document()
    document.add_paragraph("확장자만 바꾼 DOCX")
    document.save(path)

    # DOCX, PPTX, XLSX는 모두 ZIP Magic Byte를 사용한다. 선택된 파서가
    # 패키지 내부의 형식별 필수 경로까지 확인해야 위장 파일을 차단할 수 있다.
    with pytest.raises(InvalidDocumentError):
        await PptxDocumentParser().parse(path)


@pytest.mark.asyncio
async def test_txt_detects_utf16_and_preserves_line_numbers(tmp_path: Path) -> None:
    path = tmp_path / "sample.txt"
    path.write_bytes("첫 줄\n둘째 줄".encode("utf-16"))

    parsed = await TxtDocumentParser().parse(path)

    assert parsed.file_type is DocumentType.TXT
    assert parsed.document_metadata["byte_order_mark"] is True
    assert [unit.source_metadata["line_number"] for unit in parsed.units] == [1, 2]
    assert parsed.text == "첫 줄\n\n둘째 줄"


@pytest.mark.asyncio
async def test_txt_detects_utf16_without_bom(tmp_path: Path) -> None:
    path = tmp_path / "utf16-no-bom.txt"
    path.write_bytes("첫 줄\n둘째 줄".encode("utf-16-le"))

    parsed = await TxtDocumentParser().parse(path)

    assert parsed.document_metadata["encoding"].replace("-", "_").lower() == "utf_16_le"
    assert parsed.text == "첫 줄\n\n둘째 줄"


@pytest.mark.asyncio
async def test_txt_prefers_cp949_when_generic_detector_misreads_short_korean(
    tmp_path: Path,
) -> None:
    path = tmp_path / "cp949.txt"
    path.write_bytes("첫 줄\n둘째 줄".encode("cp949"))

    parsed = await TxtDocumentParser().parse(path)

    assert parsed.document_metadata["encoding"].lower() == "cp949"
    assert parsed.text == "첫 줄\n\n둘째 줄"


@pytest.mark.asyncio
async def test_txt_rejects_binary_file(tmp_path: Path) -> None:
    path = tmp_path / "binary.txt"
    path.write_bytes(b"\x89PNG\r\n\x1a\n\x00\x01")

    with pytest.raises(InvalidDocumentError):
        await TxtDocumentParser().parse(path)
