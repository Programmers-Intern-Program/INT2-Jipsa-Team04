"""다운로드 임시 확장자를 사용하는 PPTX 이미지 추출 회귀를 검증한다."""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any, cast

import pytest
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from jipsa_rag.core.document_processing import DocumentProcessingSettings
from jipsa_rag.infrastructure.document.images import download_safe_office
from jipsa_rag.infrastructure.document.images.download_safe_pptx import (
    DownloadSafePptxImageExtractor,
)
from jipsa_rag.infrastructure.document.images.models import (
    DocumentImageKind,
)
from jipsa_rag.infrastructure.document.rendering import OfficeVisualRenderResult


class _RecordingRenderer:
    """PowerPoint 호출 경로와 호출 시점의 파일 존재 여부를 기록하는 렌더러 대역."""

    def __init__(
        self,
        *,
        result: OfficeVisualRenderResult | None = None,
        error: Exception | None = None,
    ) -> None:
        self._result = result or OfficeVisualRenderResult.unavailable("unit_test")
        self._error = error
        self.pptx_paths: list[Path] = []
        self.path_exists_during_call: list[bool] = []

    async def render_pptx_visuals(
        self,
        source_path: Path,
    ) -> OfficeVisualRenderResult:
        self.pptx_paths.append(source_path)
        self.path_exists_during_call.append(source_path.is_file())

        if self._error is not None:
            raise self._error

        return self._result

    async def render_xlsx_charts(
        self,
        source_path: Path,
    ) -> OfficeVisualRenderResult:
        del source_path
        return OfficeVisualRenderResult.unavailable("unit_test")


def _settings() -> DocumentProcessingSettings:
    """그림 필터를 끄고 Office 렌더링 경로를 활성화한 결정적 설정을 생성한다."""

    return DocumentProcessingSettings(
        image_extraction_enabled=True,
        image_decorative_filter_enabled=False,
        office_rendering_enabled=True,
        ocr_enabled=False,
        ocr_gpu=False,
        ocr_gpu_required=False,
        _env_file=None,
    )


def _create_pptx(path: Path, image_path: Path) -> None:
    """삽입 그림과 차트가 함께 있는 최소 PPTX fixture를 생성한다."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image_path),
        Inches(0.5),
        Inches(0.5),
        width=Inches(3),
    )

    chart_data_factory = cast(Any, ChartData)
    chart_data = chart_data_factory()
    chart_data.categories = ["A", "B", "C"]
    cast(Any, chart_data).add_series("Series", (1, 3, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(4),
        Inches(1),
        Inches(5),
        Inches(4),
        chart_data,
    )
    presentation.save(str(path))


def _create_png(path: Path) -> None:
    """삽입 이미지와 테스트 입력으로 사용할 결정적인 PNG를 생성한다."""

    Image.new("RGB", (320, 120), "white").save(path, format="PNG")


@pytest.mark.asyncio
async def test_extracts_pptx_from_document_extension_with_source_pptx_path(
    tmp_path: Path,
) -> None:
    """``*.document`` PPTX의 ZIP 분석과 COM 호출이 같은 ``source.pptx``를 사용한다."""

    image_path = tmp_path / "sample.png"
    source_path = tmp_path / "source.pptx"
    downloaded_path = tmp_path / "downloaded.document"
    _create_png(image_path)
    _create_pptx(source_path, image_path)
    downloaded_path.write_bytes(source_path.read_bytes())
    original_bytes = downloaded_path.read_bytes()

    renderer = _RecordingRenderer(
        result=OfficeVisualRenderResult.unavailable("timeout"),
    )
    extractor = DownloadSafePptxImageExtractor(_settings(), renderer)

    extraction = await extractor.extract(downloaded_path)

    assert renderer.path_exists_during_call == [True]
    assert len(renderer.pptx_paths) == 1
    normalized_path = renderer.pptx_paths[0]
    assert normalized_path.name == "source.pptx"
    assert normalized_path.suffix == ".pptx"
    assert normalized_path != downloaded_path

    assert any(image.kind is DocumentImageKind.PPTX_PICTURE for image in extraction.images)
    assert extraction.document_metadata["office_renderer_available"] is False

    # TemporaryDirectory 수명은 extract 호출로 제한된다. timeout 부분 실패 결과를
    # 반환한 직후에도 복사본과 임시 디렉터리가 모두 제거되어야 한다.
    assert not normalized_path.exists()
    assert not normalized_path.parent.exists()

    # 안전 어댑터는 다운로더가 만든 원본을 삭제하거나 변경하지 않는다.
    assert downloaded_path.read_bytes() == original_bytes


@pytest.mark.asyncio
async def test_uses_original_pptx_path_without_copy(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """실제 ``.pptx`` 경로는 추가 복사 없이 기존 추출기로 직접 전달한다."""

    image_path = tmp_path / "sample.png"
    source_path = tmp_path / "already-normalized.pptx"
    _create_png(image_path)
    _create_pptx(source_path, image_path)

    def fail_copyfile(source: object, destination: object) -> None:
        del source, destination
        raise AssertionError("A normalized .pptx path must not be copied.")

    monkeypatch.setattr(
        shutil,
        "copyfile",
        fail_copyfile,
    )

    renderer = _RecordingRenderer()
    extractor = DownloadSafePptxImageExtractor(_settings(), renderer)

    extraction = await extractor.extract(source_path)

    assert renderer.pptx_paths == [source_path]
    assert any(image.kind is DocumentImageKind.PPTX_PICTURE for image in extraction.images)


@pytest.mark.asyncio
async def test_removes_temporary_pptx_after_renderer_exception(
    tmp_path: Path,
) -> None:
    """렌더러 예외가 상위로 전파돼도 임시 ``source.pptx``를 남기지 않는다."""

    image_path = tmp_path / "sample.png"
    source_path = tmp_path / "source.pptx"
    downloaded_path = tmp_path / "downloaded.document"
    _create_png(image_path)
    _create_pptx(source_path, image_path)
    downloaded_path.write_bytes(source_path.read_bytes())

    renderer = _RecordingRenderer(error=RuntimeError("synthetic renderer failure"))
    extractor = DownloadSafePptxImageExtractor(_settings(), renderer)

    with pytest.raises(RuntimeError, match="synthetic renderer failure"):
        await extractor.extract(downloaded_path)

    assert len(renderer.pptx_paths) == 1
    normalized_path = renderer.pptx_paths[0]
    assert not normalized_path.exists()
    assert not normalized_path.parent.exists()
    assert downloaded_path.is_file()


def test_rejects_unsupported_office_suffix() -> None:
    """공통 경로 계층은 허용되지 않은 Office 확장자를 명확히 거부한다."""

    with pytest.raises(ValueError, match="Unsupported Office document suffix"):
        download_safe_office._normalize_expected_suffix(".docx")
