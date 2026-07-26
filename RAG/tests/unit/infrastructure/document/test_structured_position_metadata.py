"""PPTX, XLSX, TXT 위치 메타데이터와 비정상 OOXML 방어를 테스트한다."""

from pathlib import Path
from typing import cast
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from openpyxl import Workbook  # type: ignore[import-untyped]
from openpyxl.worksheet.table import (  # type: ignore[import-untyped]
    Table,
    TableStyleInfo,
)
from pptx import Presentation
from pptx.util import Inches

from jipsa_rag.core.document_processing import get_document_processing_settings
from jipsa_rag.infrastructure.document.exceptions import (
    DocumentTextNotFoundError,
    EncryptedDocumentError,
    InvalidDocumentError,
)
from jipsa_rag.infrastructure.document.models import DocumentType
from jipsa_rag.infrastructure.document.parsers._common import validate_ooxml_package
from jipsa_rag.infrastructure.document.parsers.pptx import PptxDocumentParser
from jipsa_rag.infrastructure.document.parsers.txt import TxtDocumentParser
from jipsa_rag.infrastructure.document.parsers.xlsx import XlsxDocumentParser


@pytest.mark.asyncio
async def test_pptx_parser_preserves_slide_shape_geometry(tmp_path: Path) -> None:
    """도형의 슬라이드 번호, Z 순서, EMU 좌표와 정규화 위치를 생성한다."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    shape.text = "분기별 프로젝트 현황"
    file_path = tmp_path / "position.pptx"
    presentation.save(str(file_path))

    parsed = await PptxDocumentParser().parse(file_path)
    unit = next(unit for unit in parsed.units if unit.text == "분기별 프로젝트 현황")
    metadata = unit.source_metadata

    assert metadata["slide_number"] == 1
    assert metadata["shape_index"] == 1
    assert metadata["shape_z_order"] == 1
    assert metadata["shape_path"] == "1"
    assert metadata["coordinate_space"] == "slide"
    assert metadata["shape_left_emu"] == Inches(1)
    assert metadata["shape_top_emu"] == Inches(2)
    assert metadata["shape_width_emu"] == Inches(3)
    assert metadata["shape_height_emu"] == Inches(1)
    assert metadata["shape_right_emu"] == Inches(4)
    assert metadata["shape_bottom_emu"] == Inches(3)
    shape_left_ratio = cast(float, metadata["shape_left_ratio"])
    slide_width_emu = cast(int, parsed.document_metadata["slide_width_emu"])

    assert 0 < shape_left_ratio < 1
    assert slide_width_emu > 0
    assert PptxDocumentParser().parser_version == "1.1.0"


@pytest.mark.asyncio
async def test_xlsx_parser_preserves_sheet_and_cell_range(tmp_path: Path) -> None:
    """행 unit에 시트 번호, 셀 범위, 병합 셀과 표 이름을 보존한다."""

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "매출 현황"
    sheet.append(["지역", "1분기", "2분기"])
    sheet.append(["서울", 10, 20])
    sheet.append(["부산", 30, 40])
    sheet["D5"] = "합계"
    sheet.merge_cells("D5:E5")

    table = Table(displayName="SalesTable", ref="A1:C3")
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True)
    sheet.add_table(table)

    file_path = tmp_path / "position.xlsx"
    workbook.save(file_path)
    workbook.close()

    parsed = await XlsxDocumentParser().parse(file_path)
    first_row = next(unit for unit in parsed.units if unit.source_metadata["row_number"] == 1)
    merged_row = next(unit for unit in parsed.units if unit.source_metadata["row_number"] == 5)

    assert first_row.source_metadata["sheet_index"] == 1
    assert first_row.source_metadata["sheet_number"] == 1
    assert first_row.source_metadata["sheet_name"] == "매출 현황"
    assert first_row.source_metadata["start_cell"] == "A1"
    assert first_row.source_metadata["end_cell"] == "C1"
    assert first_row.source_metadata["cell_range"] == "A1:C1"
    assert first_row.source_metadata["table_names"] == ("SalesTable",)
    assert merged_row.source_metadata["merged_cell_ranges"] == ("D5:E5",)
    assert XlsxDocumentParser().parser_version == "1.1.0"


@pytest.mark.asyncio
async def test_txt_parser_preserves_line_and_character_ranges(tmp_path: Path) -> None:
    """CRLF를 LF로 정규화한 뒤 줄 번호와 exclusive 문자 범위를 계산한다."""

    file_path = tmp_path / "position.txt"
    file_path.write_bytes("가  \r\nABC\n".encode())

    parsed = await TxtDocumentParser().parse(file_path)
    first, second, trailing_empty = parsed.units

    assert first.text == "가"
    assert first.source_metadata["line_number"] == 1
    assert first.source_metadata["source_char_start"] == 0
    assert first.source_metadata["source_char_end"] == 3
    assert first.source_metadata["text_char_end"] == 1
    assert first.source_metadata["line_break_end"] == 4

    assert second.source_metadata["source_char_start"] == 4
    assert second.source_metadata["source_char_end"] == 7
    assert second.source_metadata["line_break_end"] == 8
    assert trailing_empty.source_metadata["source_char_start"] == 8
    assert trailing_empty.source_metadata["source_char_end"] == 8
    assert parsed.document_metadata["normalized_char_count"] == 8
    assert TxtDocumentParser().parser_version == "1.1.0"


@pytest.mark.asyncio
async def test_txt_parser_rejects_empty_file(tmp_path: Path) -> None:
    """0바이트 TXT는 검색 가능한 문서가 아니므로 명시적으로 거부한다."""

    file_path = tmp_path / "empty.txt"
    file_path.write_bytes(b"")

    with pytest.raises(DocumentTextNotFoundError):
        await TxtDocumentParser().parse(file_path)


@pytest.mark.asyncio
async def test_pptx_parser_rejects_corrupt_zip(tmp_path: Path) -> None:
    """PK 시그니처만 가진 손상 PPTX를 InvalidDocumentError로 변환한다."""

    file_path = tmp_path / "corrupt.pptx"
    file_path.write_bytes(b"PK\x03\x04not-a-valid-central-directory")

    with pytest.raises(InvalidDocumentError):
        await PptxDocumentParser().parse(file_path)


@pytest.mark.asyncio
async def test_xlsx_parser_rejects_ole_encrypted_document(tmp_path: Path) -> None:
    """Office 암호화 OLE 컨테이너를 일반 손상 ZIP과 구분한다."""

    file_path = tmp_path / "encrypted.xlsx"
    file_path.write_bytes(
        b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + "EncryptedPackage".encode("utf-16-le")
    )

    with pytest.raises(EncryptedDocumentError):
        await XlsxDocumentParser().parse(file_path)


@pytest.mark.asyncio
async def test_xlsx_parser_detects_ole_encryption_marker_near_file_tail(
    tmp_path: Path,
) -> None:
    """큰 OLE 파일 뒤쪽의 암호화 스트림 이름도 제한 메모리로 감지한다."""

    file_path = tmp_path / "encrypted-tail.xlsx"
    file_path.write_bytes(
        b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
        + (b"\x00" * (1024 * 1024 + 128))
        + "EncryptionInfo".encode("utf-16-le")
    )

    with pytest.raises(EncryptedDocumentError):
        await XlsxDocumentParser().parse(file_path)


def test_ooxml_validator_rejects_excessive_compression_ratio(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """압축 해제 비율 한계를 넘는 OOXML ZIP Bomb 후보를 거부한다."""

    file_path = tmp_path / "compressed.docx"
    with ZipFile(file_path, "w", compression=ZIP_DEFLATED) as package:
        package.writestr("[Content_Types].xml", "A" * 100_000)
        package.writestr("word/document.xml", "<document />")

    monkeypatch.setenv("JIPSA_RAG_OOXML_MAX_COMPRESSION_RATIO", "1.01")
    get_document_processing_settings.cache_clear()
    try:
        with pytest.raises(InvalidDocumentError):
            validate_ooxml_package(
                file_path,
                file_type=DocumentType.DOCX,
                required_members=frozenset(
                    {
                        "[Content_Types].xml",
                        "word/document.xml",
                    }
                ),
            )
    finally:
        get_document_processing_settings.cache_clear()
