"""PPTX 슬라이드의 텍스트와 도형 위치 메타데이터를 추출한다.

슬라이드 번호만 저장하면 사용자가 본 화면에서 근거가 어느 도형에 있었는지
재현하기 어렵다. 이 파서는 슬라이드의 읽기 순서를 유지하면서 제목, 일반 텍스트
도형, 표, 그룹 내부 도형과 발표자 노트를 독립적인 ``ParsedDocumentUnit``으로
변환하고, 텍스트가 있는 시각 도형에는 다음 위치 계약을 함께 기록한다.

- 1부터 시작하는 슬라이드 번호와 최상위 Z 순서
- python-pptx의 안정적인 ``shape_id``와 중첩 ``shape_path``
- EMU 단위의 left/top/width/height
- 슬라이드 크기에 대한 0~1 정규화 위치와 크기
- 그룹 내부 도형인지 구분하는 좌표 공간과 중첩 깊이

EMU(English Metric Unit)는 Office Open XML이 사용하는 정수 좌표 단위다. 픽셀은
렌더링 DPI에 따라 달라지므로 원본 위치 계약에는 EMU와 정규화 비율을 사용한다.
차트 데이터, SmartArt 의미 구조와 이미지 OCR은 현재 범위에 포함하지 않는다.
"""

import asyncio
from collections.abc import Iterable
from pathlib import Path
from typing import Final, Protocol, cast

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as PptxPresentation
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.slide import Slide

from jipsa_rag.infrastructure.document.exceptions import (
    DocumentParserError,
    DocumentReadError,
    DocumentTextExtractionError,
    DocumentTextNotFoundError,
    InvalidDocumentError,
)
from jipsa_rag.infrastructure.document.models import (
    DocumentType,
    ParsedDocument,
    ParsedDocumentUnit,
    SourceMetadataValue,
)
from jipsa_rag.infrastructure.document.parsers._common import (
    normalize_inline_text,
    normalize_text,
    validate_ooxml_package,
)

_PPTX_PARSER_TYPE: Final[str] = "PPTX_TEXT"

# 1.1.0부터 도형의 EMU 좌표, 정규화 비율, shape_id와 좌표 공간을 보존한다.
# Parser_Version은 Local RAG 문서 유일성, 결정적 Chunk ID와 Qdrant payload에
# 포함되므로 기존 1.0.0 색인은 자동으로 재파싱·재색인 대상이 된다.
_PPTX_PARSER_VERSION: Final[str] = "1.1.0"

_PPTX_REQUIRED_MEMBERS: Final[frozenset[str]] = frozenset(
    {
        "[Content_Types].xml",
        "ppt/presentation.xml",
    }
)


class _ShapeWithText(Protocol):
    """``has_text_frame``이 참인 도형에서 사용하는 최소 텍스트 계약."""

    @property
    def text(self) -> str:
        """도형 텍스트 프레임의 전체 문자열을 반환한다."""

        ...


class PptxDocumentParser:
    """``python-pptx``로 슬라이드 구조와 도형 위치를 보존해 파싱한다."""

    @property
    def file_type(self) -> DocumentType:
        """이 파서가 처리하는 문서 형식을 반환한다."""

        return DocumentType.PPTX

    @property
    def parser_type(self) -> str:
        """Local RAG DB에 저장할 파서 종류를 반환한다."""

        return _PPTX_PARSER_TYPE

    @property
    def parser_version(self) -> str:
        """재파싱과 결정적 Chunk ID에 사용하는 파서 버전을 반환한다."""

        return _PPTX_PARSER_VERSION

    async def parse(self, file_path: Path) -> ParsedDocument:
        """동기식 ZIP/XML 파싱을 작업 스레드에서 실행한다."""

        return await asyncio.to_thread(self._parse_sync, file_path)

    def _parse_sync(self, file_path: Path) -> ParsedDocument:
        """PPTX 패키지를 검증하고 슬라이드 순서대로 unit을 생성한다."""

        validate_ooxml_package(
            file_path,
            file_type=self.file_type,
            required_members=_PPTX_REQUIRED_MEMBERS,
        )

        try:
            presentation = Presentation(str(file_path))
            # python-pptx 타입 계약상 슬라이드 높이는 None일 수 있다.
            # 실제 PPTX 파싱에 필요한 크기가 누락되면 위치 비율을 임의로 만들지
            # 않고 손상 문서로 분류한다.
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height

            if slide_width is None or slide_height is None:
                raise InvalidDocumentError(self.file_type)

            slide_width_emu = int(slide_width)
            slide_height_emu = int(slide_height)
            units = self._parse_presentation(
                presentation,
                slide_width_emu=slide_width_emu,
                slide_height_emu=slide_height_emu,
            )
        except DocumentParserError:
            raise
        except OSError as error:
            raise DocumentReadError(file_path) from error
        except Exception as error:
            raise InvalidDocumentError(self.file_type) from error

        parsed_document = ParsedDocument(
            file_type=self.file_type,
            units=tuple(units),
            document_metadata={
                "slide_count": len(presentation.slides),
                "source_unit_count": len(units),
                "slide_width_emu": slide_width_emu,
                "slide_height_emu": slide_height_emu,
            },
        )

        if parsed_document.text_unit_count == 0:
            raise DocumentTextNotFoundError(self.file_type)

        return parsed_document

    def _parse_presentation(
        self,
        presentation: PptxPresentation,
        *,
        slide_width_emu: int,
        slide_height_emu: int,
    ) -> list[ParsedDocumentUnit]:
        """모든 슬라이드를 1부터 시작하는 번호 순서로 파싱한다."""

        units: list[ParsedDocumentUnit] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            title_shape = slide.shapes.title
            title_shape_id = title_shape.shape_id if title_shape is not None else None

            for shape_index, shape in enumerate(slide.shapes, start=1):
                try:
                    units.extend(
                        self._parse_shape(
                            shape,
                            slide_number=slide_number,
                            shape_index=shape_index,
                            shape_path=str(shape_index),
                            is_title=(
                                title_shape_id is not None and shape.shape_id == title_shape_id
                            ),
                            slide_width_emu=slide_width_emu,
                            slide_height_emu=slide_height_emu,
                            coordinate_space="slide",
                        )
                    )
                except DocumentParserError:
                    raise
                except Exception as error:
                    raise DocumentTextExtractionError(
                        file_type=self.file_type,
                        source_metadata={
                            "slide_number": slide_number,
                            "shape_index": shape_index,
                        },
                    ) from error

            notes_text = self._extract_notes_text(slide)
            if notes_text:
                units.append(
                    ParsedDocumentUnit(
                        text=notes_text,
                        source_metadata={
                            "unit_type": "speaker_notes",
                            "location_kind": "pptx_speaker_notes",
                            "slide_number": slide_number,
                            "notes_index": 1,
                        },
                    )
                )

        return units

    def _parse_shape(
        self,
        shape: BaseShape,
        *,
        slide_number: int,
        shape_index: int,
        shape_path: str,
        is_title: bool,
        slide_width_emu: int,
        slide_height_emu: int,
        coordinate_space: str,
    ) -> Iterable[ParsedDocumentUnit]:
        """단일 도형 또는 그룹 도형을 위치 메타데이터가 있는 unit으로 변환한다.

        ``shape_index``는 슬라이드 최상위 도형의 Z 순서이며 그룹 자식도 부모의
        값을 공유한다. 실제 그룹 내부 순서는 ``shape_path``의 점 구분 경로로
        표현한다. 그룹 자식 좌표는 그룹 변환의 영향을 받으므로
        ``coordinate_space=group``을 명시하여 소비자가 슬라이드 절대 좌표로
        오해하지 않게 한다.
        """

        metadata: dict[str, SourceMetadataValue] = {
            "location_kind": "pptx_shape",
            "slide_number": slide_number,
            "shape_index": shape_index,
            "shape_z_order": shape_index,
            "shape_id": int(shape.shape_id),
            "shape_path": shape_path,
            "shape_depth": shape_path.count(".") + 1,
            "shape_name": shape.name,
            "shape_type": int(shape.shape_type),
            "shape_type_name": getattr(
                shape.shape_type,
                "name",
                str(shape.shape_type),
            ),
            "coordinate_space": coordinate_space,
            **_shape_geometry_metadata(
                shape,
                slide_width_emu=slide_width_emu,
                slide_height_emu=slide_height_emu,
            ),
        }

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and isinstance(shape, GroupShape):
            grouped_units: list[ParsedDocumentUnit] = []

            for child_index, child_shape in enumerate(shape.shapes, start=1):
                grouped_units.extend(
                    self._parse_shape(
                        child_shape,
                        slide_number=slide_number,
                        shape_index=shape_index,
                        shape_path=f"{shape_path}.{child_index}",
                        is_title=False,
                        slide_width_emu=slide_width_emu,
                        slide_height_emu=slide_height_emu,
                        coordinate_space="group",
                    )
                )

            return grouped_units

        if shape.has_table:
            table = cast(GraphicFrame, shape).table
            rows: list[str] = []

            for row in table.rows:
                row_text = "\t".join(
                    normalize_inline_text(cell.text) for cell in row.cells
                ).rstrip()
                rows.append(row_text)

            return (
                ParsedDocumentUnit(
                    text=normalize_text("\n".join(rows)),
                    source_metadata={
                        **metadata,
                        "unit_type": "table",
                        "row_count": len(table.rows),
                        "column_count": len(table.columns),
                    },
                ),
            )

        if shape.has_text_frame:
            text = normalize_text(cast(_ShapeWithText, shape).text)
            unit_type = "title" if is_title else "shape_text"

            return (
                ParsedDocumentUnit(
                    text=text,
                    source_metadata={
                        **metadata,
                        "unit_type": unit_type,
                        "is_title": is_title,
                    },
                ),
            )

        return ()

    @staticmethod
    def _extract_notes_text(slide: Slide) -> str:
        """발표자 노트 본문을 반환하고 비표준 노트 관계는 빈 값으로 처리한다."""

        if not slide.has_notes_slide:
            return ""

        try:
            notes_text_frame = slide.notes_slide.notes_text_frame
        except (AttributeError, KeyError):
            return ""

        if notes_text_frame is None:
            return ""

        return normalize_text(notes_text_frame.text)


def _shape_geometry_metadata(
    shape: BaseShape,
    *,
    slide_width_emu: int,
    slide_height_emu: int,
) -> dict[str, SourceMetadataValue]:
    """도형의 EMU 좌표와 슬라이드 기준 정규화 위치를 생성한다.

    정규화 비율은 서로 다른 슬라이드 크기의 PPTX를 동일한 방식으로 표시하기 위한
    보조 정보다. 원본 정밀도를 잃지 않도록 EMU 정수도 반드시 함께 보존한다.
    """

    left_emu = int(shape.left)
    top_emu = int(shape.top)
    width_emu = int(shape.width)
    height_emu = int(shape.height)

    metadata: dict[str, SourceMetadataValue] = {
        "shape_left_emu": left_emu,
        "shape_top_emu": top_emu,
        "shape_width_emu": width_emu,
        "shape_height_emu": height_emu,
        "shape_right_emu": left_emu + width_emu,
        "shape_bottom_emu": top_emu + height_emu,
        "shape_left_ratio": _safe_ratio(left_emu, slide_width_emu),
        "shape_top_ratio": _safe_ratio(top_emu, slide_height_emu),
        "shape_width_ratio": _safe_ratio(width_emu, slide_width_emu),
        "shape_height_ratio": _safe_ratio(height_emu, slide_height_emu),
    }

    rotation = getattr(shape, "rotation", None)
    if isinstance(rotation, int | float):
        metadata["shape_rotation_degrees"] = float(rotation)

    return metadata


def _safe_ratio(value: int, denominator: int) -> float:
    """EMU 값을 0으로 나누지 않고 결정적인 8자리 비율로 변환한다."""

    if denominator <= 0:
        return 0.0
    return round(value / denominator, 8)
